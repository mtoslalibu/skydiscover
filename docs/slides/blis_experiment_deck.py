#!/usr/bin/env python3
"""BLIS Router Experiment slide deck — 2 versions per topic, 20 slides total."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG     = RGBColor(0x16, 0x16, 0x2B)
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)
ACCENT_TEAL = RGBColor(0x00, 0xB4, 0xA0)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x42)
ACCENT_PURPLE = RGBColor(0x9B, 0x59, 0xB6)
LIGHT_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY    = RGBColor(0x88, 0x88, 0x99)
SUBTLE_GRAY = RGBColor(0x44, 0x44, 0x55)
GREEN       = RGBColor(0x2E, 0xCC, 0x71)
RED         = RGBColor(0xE7, 0x4C, 0x3C)
CARD_BG     = RGBColor(0x22, 0x22, 0x3A)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


# ── Helpers ──
def set_bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def T(slide, left, top, width, height, txt, size=13, bold=False,
      color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align
    return tb


def MT(slide, left, top, width, height, lines, ds=12, dc=LIGHT_GRAY):
    """Multi-line textbox. lines = list of str or (text,size,bold,color)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            t, s, b, c = item, ds, False, dc
        else:
            t = item[0]
            s = item[1] if len(item) > 1 else ds
            b = item[2] if len(item) > 2 else False
            c = item[3] if len(item) > 3 else dc
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = t
        p.font.size = Pt(s)
        p.font.bold = b
        p.font.color.rgb = c
        p.font.name = "Calibri"
        p.space_after = Pt(3)
    return tb


def R(slide, left, top, width, height, fill=CARD_BG, line=None,
      shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    return s


def BAR(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def hdr(slide, title, sub="", sc=ACCENT_TEAL):
    T(slide, Inches(0.6), Inches(0.2), Inches(12.2), Inches(0.65),
      title, size=25, bold=True)
    if sub:
        T(slide, Inches(0.6), Inches(0.82), Inches(11.5), Inches(0.38),
          sub, size=13, color=sc)


def ftr(slide, msg):
    T(slide, Inches(0.6), Inches(7.12), Inches(12.1), Inches(0.28),
      msg, size=9, color=MID_GRAY)


def vtag(slide, v):
    T(slide, Inches(11.8), Inches(0.12), Inches(1.2), Inches(0.32),
      f"v{v}", size=10, bold=True, color=MID_GRAY, align=PP_ALIGN.RIGHT)


def card(slide, left, top, w, h, accent, title, lines):
    R(slide, left, top, w, h, fill=CARD_BG, line=accent)
    BAR(slide, left, top, w, Inches(0.055), accent)
    T(slide, left+Inches(0.15), top+Inches(0.09), w-Inches(0.3), Inches(0.34),
      title, size=12, bold=True, color=accent)
    MT(slide, left+Inches(0.15), top+Inches(0.48), w-Inches(0.3),
       h-Inches(0.58), lines, ds=11, dc=LIGHT_GRAY)


# ════════════════════════════════════════════════════════
# SLIDE 1A — Title (impact-first)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); vtag(sl, "A")
hdr(sl, "BLIS Router Optimization — Experiment Results",
    "SkyDiscover × OpenEvolve  •  50 iterations  •  March 2026")

R(sl, Inches(0.6), Inches(1.4), Inches(5.6), Inches(2.9), fill=CARD_BG, line=GREEN)
T(sl, Inches(1.0), Inches(1.55), Inches(5.0), Inches(0.9),
  "+11.5%", size=52, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
T(sl, Inches(0.8), Inches(2.5), Inches(5.2), Inches(0.4),
  "improvement over best hand-tuned baseline (1:1)", size=13,
  color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
MT(sl, Inches(0.8), Inches(3.0), Inches(5.2), Inches(1.1), [
    ("  Found at iteration 16 of 50 — just ~10 minutes", 12, False, MID_GRAY),
    ("  0 build errors across all 50 iterations", 12, False, MID_GRAY),
    ("  Beats the human-crafted oracle by +0.6%", 12, True, GREEN),
], ds=12)

R(sl, Inches(6.5), Inches(1.4), Inches(6.3), Inches(2.9), fill=CARD_BG)
T(sl, Inches(6.7), Inches(1.5), Inches(6.0), Inches(0.38),
  "What this deck covers", size=13, bold=True, color=WHITE)
MT(sl, Inches(6.7), Inches(1.92), Inches(5.9), Inches(2.25), [
    ("→  How SkyDiscover + OpenEvolve works (workflow)", 12, False, LIGHT_GRAY),
    ("→  LLMs used as mutation engine (Claude Sonnet + Opus)", 12, False, LIGHT_GRAY),
    ("→  Seeds: what they are, why two", 12, False, LIGHT_GRAY),
    ("→  Mutator: what the LLM actually changed in the code", 12, False, LIGHT_GRAY),
    ("→  Prompt: what context we gave the LLM", 12, False, LIGHT_GRAY),
    ("→  Two workloads: general vs prefix-heavy traffic", 12, False, LIGHT_GRAY),
    ("→  Baselines: 6 algorithms we compared against", 12, False, LIGHT_GRAY),
    ("→  Results: per-workload breakdown with all baselines", 12, False, LIGHT_GRAY),
    ("→  Sim2Real: packaging the algorithm for real deployment", 12, False, LIGHT_GRAY),
], ds=12)

ftr(sl, "Experiment: 260312_50i_openevolve_v2wl  •  qwen_7b eval  •  seeds 42+456  •  4×H100 cluster simulation")

# ════════════════════════════════════════════════════════
# SLIDE 1B — Title (question-driven)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); vtag(sl, "B")
T(sl, Inches(0.6), Inches(0.75), Inches(12.2), Inches(0.9),
  "Can an AI discover a better LLM routing algorithm\nthan a human expert?",
  size=28, bold=True, color=WHITE)
T(sl, Inches(0.6), Inches(1.75), Inches(12), Inches(0.45),
  "We gave SkyDiscover a routing function and an evaluator. It found +11.5% improvement in 10 minutes.",
  size=16, color=LIGHT_GRAY)
T(sl, Inches(0.6), Inches(2.35), Inches(12), Inches(0.5),
  "Yes — and it independently rediscovered the same core idea as the human-crafted oracle.",
  size=18, bold=True, color=GREEN)

cw = Inches(3.8); cy = Inches(3.05)
for i, (title, accent, blines) in enumerate([
    ("The Problem", ACCENT_BLUE, [
        ("LLM clusters need smart routing to minimize latency", 12, False, LIGHT_GRAY),
        ("", 5),
        ("One prefix group can dominate → hotspot on", 12, False, LIGHT_GRAY),
        ("one server → queue builds up → slow responses", 12, False, LIGHT_GRAY),
        ("", 5),
        ("Hand-tuning routing weights is slow & hard", 12, False, LIGHT_GRAY),
    ]),
    ("Our Approach", ACCENT_TEAL, [
        ("SkyDiscover + OpenEvolve: LLM evolves Go code", 12, False, LIGHT_GRAY),
        ("", 5),
        ("Each iteration: mutate → build → simulate → score", 12, False, LIGHT_GRAY),
        ("", 5),
        ("50 iters, 2 seeds × 2 workloads per score", 12, False, LIGHT_GRAY),
        ("Claude Sonnet (70%) + Opus (30%) as mutators", 12, False, LIGHT_GRAY),
    ]),
    ("The Result", GREEN, [
        ("+11.5% over best hand-tuned 1:1 baseline", 12, True, GREEN),
        ("", 5),
        ("+11.4% on prefix-heavy workload", 12, False, LIGHT_GRAY),
        ("+0.3%  on general-traffic workload", 12, False, LIGHT_GRAY),
        ("", 5),
        ("0 build errors  •  converged at iter 16", 12, False, MID_GRAY),
        ("Beats human oracle v14 by +0.6%", 12, True, GREEN),
    ]),
]):
    x = Inches(0.5) + i * (cw + Inches(0.45))
    R(sl, x, cy, cw, Inches(3.85), fill=CARD_BG, line=accent)
    BAR(sl, x, cy, cw, Inches(0.055), accent)
    T(sl, x+Inches(0.15), cy+Inches(0.1), cw-Inches(0.3), Inches(0.38),
      title, size=14, bold=True, color=accent)
    MT(sl, x+Inches(0.15), cy+Inches(0.55), cw-Inches(0.3), Inches(3.1), blines)

ftr(sl, "Experiment: 260312_50i_openevolve_v2wl  •  March 2026")


# ════════════════════════════════════════════════════════
# SLIDE 2A — Workflow (step-by-step)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); vtag(sl, "A")
hdr(sl, "How SkyDiscover + OpenEvolve Works",
    "Automated loop: LLM mutates code → simulator scores it → repeat 50×")

# Claude orchestrator steps (left)
T(sl, Inches(0.6), Inches(1.42), Inches(5.8), Inches(0.32),
  "Claude (Experiment Orchestrator)", size=13, bold=True, color=ACCENT_BLUE)
orch_steps = [
    ("1. Setup",   "Reads config, sets env vars, creates output dir", ACCENT_BLUE),
    ("2. Launch",  "Runs skydiscover-run CLI → hands off to OpenEvolve", ACCENT_TEAL),
    ("3. Monitor", "Checks logs every 2 min: best score, errors, progress", ACCENT_ORANGE),
    ("4. Analyze", "Runs 4 scripts: compare, plot, effort, diffs", ACCENT_PURPLE),
    ("5. Report",  "Writes analysis.md with tables & key findings", GREEN),
]
for i, (name, desc, color) in enumerate(orch_steps):
    y = Inches(1.82) + i * Inches(0.92)
    BAR(sl, Inches(0.6), y, Inches(0.1), Inches(0.73), color)
    R(sl, Inches(0.82), y, Inches(5.2), Inches(0.73), fill=CARD_BG)
    T(sl, Inches(1.02), y+Inches(0.08), Inches(1.55), Inches(0.28),
      name, size=11, bold=True, color=color)
    T(sl, Inches(2.6), y+Inches(0.08), Inches(3.3), Inches(0.58),
      desc, size=10, color=LIGHT_GRAY)
    if i < 4:
        T(sl, Inches(1.5), y+Inches(0.75), Inches(0.4), Inches(0.2),
          "▼", size=10, color=MID_GRAY, align=PP_ALIGN.CENTER)

# OpenEvolve inner loop (right)
T(sl, Inches(6.9), Inches(1.42), Inches(6.0), Inches(0.32),
  "OpenEvolve Inner Loop (runs autonomously, 50×)", size=13, bold=True, color=ACCENT_TEAL)
loop = [
    ("LLM mutates\nGo code", ACCENT_BLUE),
    ("Build &\nvalidate", ACCENT_TEAL),
    ("4 sims\n(2×seeds\n×2 WLs)", ACCENT_ORANGE),
    ("Score &\narchive", ACCENT_PURPLE),
    ("Next\niteration", GREEN),
]
lx = Inches(6.8); bw = Inches(1.95); bh = Inches(1.95); aw = Inches(0.35)
for i, (lbl, color) in enumerate(loop):
    R(sl, lx, Inches(1.82), bw, bh, fill=CARD_BG, line=color)
    num = sl.shapes.add_shape(MSO_SHAPE.OVAL, lx+Inches(0.1), Inches(1.9),
                              Inches(0.4), Inches(0.4))
    num.fill.solid(); num.fill.fore_color.rgb = color; num.line.fill.background()
    tf = num.text_frame; p = tf.paragraphs[0]
    p.text = str(i+1); p.font.size = Pt(14); p.font.bold = True
    p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    T(sl, lx+Inches(0.08), Inches(2.42), bw-Inches(0.16), Inches(1.25),
      lbl, size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
    if i < 4:
        T(sl, lx+bw, Inches(2.58), aw, Inches(0.3),
          "→", size=16, color=MID_GRAY, align=PP_ALIGN.CENTER)
    lx += bw + aw

# Bottom stats bar
R(sl, Inches(0.6), Inches(6.65), Inches(12.1), Inches(0.45), fill=CARD_BG)
T(sl, Inches(0.8), Inches(6.72), Inches(11.5), Inches(0.32),
  "50 iterations  •  ~38s/iter  •  31.8 min total  •  0 build errors  •  Best at iter 16 (32% of budget)",
  size=11, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)
ftr(sl, "OpenEvolve: 5 MAP-Elites islands, ring migration, population=40, 57.5% score diversity")

# ════════════════════════════════════════════════════════
# SLIDE 2B — Workflow (inputs → engine → outputs)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); vtag(sl, "B")
hdr(sl, "How SkyDiscover + OpenEvolve Works",
    "You provide three files. SkyDiscover does the rest automatically.")

boxy = Inches(1.42); boxh = Inches(4.9)
# YOU PROVIDE
R(sl, Inches(0.5), boxy, Inches(3.55), boxh, fill=CARD_BG, line=ACCENT_BLUE)
T(sl, Inches(0.7), boxy+Inches(0.12), Inches(3.2), Inches(0.36),
  "YOU PROVIDE", size=14, bold=True, color=ACCENT_BLUE)
MT(sl, Inches(0.7), boxy+Inches(0.58), Inches(3.2), Inches(4.1), [
    ("initial_program.go", 13, True, WHITE),
    ("  Routing function with EVOLVE-BLOCK", 11, False, MID_GRAY),
    ("  markers — LLM may only edit inside", 11, False, MID_GRAY),
    ("", 7),
    ("evaluator.py", 13, True, WHITE),
    ("  Runs BLIS simulator, returns score", 11, False, MID_GRAY),
    ("  score = % latency improvement vs baseline", 11, False, MID_GRAY),
    ("", 7),
    ("config.yaml", 13, True, WHITE),
    ("  LLMs to use (Sonnet 70% + Opus 30%)", 11, False, MID_GRAY),
    ("  iterations=50, seeds, system prompt", 11, False, MID_GRAY),
])

T(sl, Inches(4.2), boxy+Inches(2.15), Inches(0.9), Inches(0.55),
  "→", size=32, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

# ENGINE
R(sl, Inches(5.2), boxy, Inches(3.4), boxh, fill=CARD_BG, line=ACCENT_TEAL)
T(sl, Inches(5.4), boxy+Inches(0.12), Inches(3.1), Inches(0.36),
  "SKYDISCOVER ENGINE", size=14, bold=True, color=ACCENT_TEAL)
MT(sl, Inches(5.4), boxy+Inches(0.58), Inches(3.1), Inches(4.1), [
    ("1. LLM generates a new routing variant", 12, False, LIGHT_GRAY),
    ("   (only modifies the EVOLVE-BLOCK)", 10, False, MID_GRAY),
    ("", 6),
    ("2. Builds & validates Go code", 12, False, LIGHT_GRAY),
    ("", 6),
    ("3. Runs 4 simulations:", 12, False, LIGHT_GRAY),
    ("   seed 42 × glia_40qps", 10, False, MID_GRAY),
    ("   seed 42 × glia_prefix_heavy", 10, False, MID_GRAY),
    ("   seed 456 × glia_40qps", 10, False, MID_GRAY),
    ("   seed 456 × glia_prefix_heavy", 10, False, MID_GRAY),
    ("", 6),
    ("4. Archives program + score", 12, False, LIGHT_GRAY),
    ("5. Repeats 50 times", 12, False, LIGHT_GRAY),
])

T(sl, Inches(8.75), boxy+Inches(2.15), Inches(0.9), Inches(0.55),
  "→", size=32, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)

# YOU GET
R(sl, Inches(9.75), boxy, Inches(3.1), boxh, fill=CARD_BG, line=ACCENT_ORANGE)
T(sl, Inches(9.95), boxy+Inches(0.12), Inches(2.8), Inches(0.36),
  "YOU GET", size=14, bold=True, color=ACCENT_ORANGE)
MT(sl, Inches(9.95), boxy+Inches(0.58), Inches(2.8), Inches(4.1), [
    ("best_program.go", 13, True, WHITE),
    ("  Drop-in replacement for routing.go", 11, False, MID_GRAY),
    ("  AI-discovered routing algorithm", 11, False, MID_GRAY),
    ("", 7),
    ("best_program_info.json", 13, True, WHITE),
    ("  Score, latencies, iteration, parent", 11, False, MID_GRAY),
    ("", 7),
    ("analysis.md + 7 plots", 13, True, WHITE),
    ("  Tables, convergence charts, diffs", 11, False, MID_GRAY),
    ("", 7),
    ("sim2real package", 13, True, WHITE),
    ("  Workloads + baselines + repro script", 11, False, MID_GRAY),
])

ftr(sl, "Total runtime: 31.8 min  •  50 iterations  •  38s avg/iter  •  best found at iteration 16")


# ════════════════════════════════════════════════════════
# SLIDE 3A — LLMs (detailed cards)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); vtag(sl, "A")
hdr(sl, "LLM Configuration: Two Models, Weighted Sampling",
    "Claude Sonnet 4-5 (70%) + Claude Opus 4-6 (30%) — used only for code mutation")

cw = Inches(5.55); ch = Inches(5.2); cy = Inches(1.42)
# Sonnet
R(sl, Inches(0.6), cy, cw, ch, fill=CARD_BG, line=ACCENT_BLUE)
BAR(sl, Inches(0.6), cy, cw, Inches(0.06), ACCENT_BLUE)
T(sl, Inches(0.82), cy+Inches(0.1), cw-Inches(0.3), Inches(0.42),
  "Claude Sonnet 4-5", size=20, bold=True, color=ACCENT_BLUE)
T(sl, Inches(0.82), cy+Inches(0.58), cw-Inches(0.3), Inches(0.32),
  "70% of all mutations", size=14, bold=True, color=WHITE)
MT(sl, Inches(0.82), cy+Inches(1.0), cw-Inches(0.3), Inches(4.0), [
    ("Why 70%?", 12, True, ACCENT_BLUE),
    ("  Fast inference → shorter iteration cycles", 11, False, LIGHT_GRAY),
    ("  High code quality at lower cost per call", 11, False, LIGHT_GRAY),
    ("  Great at incremental refinement of existing code", 11, False, LIGHT_GRAY),
    ("", 7),
    ("Model ID: aws/claude-sonnet-4-5", 11, True, WHITE),
    ("Temperature: 1.0  (diverse mutations)", 11, False, MID_GRAY),
    ("Max tokens: 32,000  (full Go code blocks)", 11, False, MID_GRAY),
    ("Timeout: 120s  •  API: IBM LiteLLM proxy", 11, False, MID_GRAY),
    ("", 7),
    ("Role in OpenEvolve:", 12, True, ACCENT_BLUE),
    ("  Primary workhorse — exploitation mode", 11, False, LIGHT_GRAY),
    ("  Refines existing good solutions in population", 11, False, LIGHT_GRAY),
])
# Opus
R(sl, Inches(7.0), cy, cw, ch, fill=CARD_BG, line=ACCENT_PURPLE)
BAR(sl, Inches(7.0), cy, cw, Inches(0.06), ACCENT_PURPLE)
T(sl, Inches(7.22), cy+Inches(0.1), cw-Inches(0.3), Inches(0.42),
  "Claude Opus 4-6", size=20, bold=True, color=ACCENT_PURPLE)
T(sl, Inches(7.22), cy+Inches(0.58), cw-Inches(0.3), Inches(0.32),
  "30% of all mutations", size=14, bold=True, color=WHITE)
MT(sl, Inches(7.22), cy+Inches(1.0), cw-Inches(0.3), Inches(4.0), [
    ("Why 30%?", 12, True, ACCENT_PURPLE),
    ("  Most capable Claude model → bigger conceptual leaps", 11, False, LIGHT_GRAY),
    ("  Better at novel algorithmic ideas & restructuring", 11, False, LIGHT_GRAY),
    ("  Slower but worth it for exploration diversity", 11, False, LIGHT_GRAY),
    ("", 7),
    ("Model ID: aws/claude-opus-4-6", 11, True, WHITE),
    ("Temperature: 1.0  (encourages creative mutations)", 11, False, MID_GRAY),
    ("Max tokens: 32,000", 11, False, MID_GRAY),
    ("", 7),
    ("Role in OpenEvolve:", 12, True, ACCENT_PURPLE),
    ("  Exploration mode — discovers new paradigms", 11, False, LIGHT_GRAY),
    ("  Finds structurally different approaches", 11, False, LIGHT_GRAY),
    ("", 7),
    ("Key: these LLMs write routing code.", 12, True, WHITE),
    ("Evaluation uses qwen_7b — a different model.", 11, False, MID_GRAY),
])
ftr(sl, "Note: LLMs = mutation engine only. Evaluation runs qwen2.5-7b-instruct on 4×H100 for fair, reproducible scoring.")

# ════════════════════════════════════════════════════════
# SLIDE 3B — LLMs (simple visual)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); vtag(sl, "B")
hdr(sl, "LLM Configuration",
    "Two Claude models working together: one fast, one creative")

R(sl, Inches(0.6), Inches(1.42), Inches(12.1), Inches(1.05), fill=CARD_BG)
T(sl, Inches(0.8), Inches(1.52), Inches(11.5), Inches(0.28),
  "Think of it like two engineers:", size=14, bold=True, color=WHITE)
T(sl, Inches(0.8), Inches(1.84), Inches(11.5), Inches(0.5),
  "Sonnet (fast, reliable) generates 7 out of 10 mutations — great for small improvements.    "
  "Opus (powerful) generates 3 out of 10 — better at brand-new ideas.",
  size=13, color=LIGHT_GRAY)

my = Inches(2.65); mh = Inches(3.5)
# Sonnet block
R(sl, Inches(0.6), my, Inches(5.85), mh, fill=CARD_BG, line=ACCENT_BLUE)
# weight bar visual
wb = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
     Inches(0.6), my+mh-Inches(0.38), Inches(5.85*0.7), Inches(0.38))
wb.fill.solid(); wb.fill.fore_color.rgb = ACCENT_BLUE; wb.line.fill.background()
T(sl, Inches(0.65), my+mh-Inches(0.38), Inches(3.0), Inches(0.38),
  "  70% of mutations", size=11, bold=True, color=WHITE)
T(sl, Inches(0.82), my+Inches(0.12), Inches(5.3), Inches(0.42),
  "Claude Sonnet 4-5", size=20, bold=True, color=ACCENT_BLUE)
T(sl, Inches(0.82), my+Inches(0.6), Inches(5.3), Inches(0.3),
  "The workhorse — fast, high-quality code edits", size=13, color=WHITE)
MT(sl, Inches(0.82), my+Inches(1.0), Inches(5.3), Inches(2.0), [
    ("✓  Short iteration cycles (model is fast)", 12, False, GREEN),
    ("✓  Reliable Go code, few syntax errors", 12, False, GREEN),
    ("✓  Great at refining existing good solutions", 12, False, GREEN),
    ("→  Temperature 1.0 for diverse outputs", 11, False, MID_GRAY),
])

# Opus block
R(sl, Inches(7.0), my, Inches(5.85), mh, fill=CARD_BG, line=ACCENT_PURPLE)
wb2 = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
      Inches(7.0), my+mh-Inches(0.38), Inches(5.85*0.3), Inches(0.38))
wb2.fill.solid(); wb2.fill.fore_color.rgb = ACCENT_PURPLE; wb2.line.fill.background()
T(sl, Inches(7.05), my+mh-Inches(0.38), Inches(1.8), Inches(0.38),
  "  30%", size=11, bold=True, color=WHITE)
T(sl, Inches(7.22), my+Inches(0.12), Inches(5.3), Inches(0.42),
  "Claude Opus 4-6", size=20, bold=True, color=ACCENT_PURPLE)
T(sl, Inches(7.22), my+Inches(0.6), Inches(5.3), Inches(0.3),
  "The creative — discovers big breakthroughs", size=13, color=WHITE)
MT(sl, Inches(7.22), my+Inches(1.0), Inches(5.3), Inches(2.0), [
    ("✓  Most powerful Claude model available", 12, False, GREEN),
    ("✓  Better at entirely new algorithm structures", 12, False, GREEN),
    ("✓  Finds what Sonnet misses", 12, False, GREEN),
    ("→  Worth the slower speed for exploration", 11, False, MID_GRAY),
])
ftr(sl, "Important: these LLMs mutate the routing code. Scoring uses qwen_7b — a separate inference model.")


# ════════════════════════════════════════════════════════
# SLIDE 4A — Seeds (technical)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); vtag(sl, "A")
hdr(sl, "Random Seeds: Robustness Through Diversity",
    "Seeds 42 and 456 — every evaluation runs both, result is the average")

R(sl, Inches(0.6), Inches(1.42), Inches(5.6), Inches(5.2), fill=CARD_BG, line=ACCENT_ORANGE)
T(sl, Inches(0.82), Inches(1.55), Inches(5.2), Inches(0.38),
  "What a seed controls", size=15, bold=True, color=ACCENT_ORANGE)
MT(sl, Inches(0.82), Inches(2.0), Inches(5.2), Inches(4.4), [
    ("A seed initializes the traffic simulator's random number generator.", 12, False, LIGHT_GRAY),
    ("", 7),
    ("It determines:", 12, True, WHITE),
    ("  •  Exact arrival times of each request", 11, False, LIGHT_GRAY),
    ("  •  Which prefix group each request uses", 11, False, LIGHT_GRAY),
    ("  •  Token lengths (drawn from distributions)", 11, False, LIGHT_GRAY),
    ("  •  Burst timing (gamma-distributed)", 11, False, LIGHT_GRAY),
    ("", 7),
    ("Same algorithm + same seed = identical result", 12, True, WHITE),
    ("Same algorithm + different seed = different traffic pattern", 12, False, LIGHT_GRAY),
    ("", 7),
    ("Using two seeds prevents overfitting:", 12, True, WHITE),
    ("  An algorithm that only works for one specific", 11, False, LIGHT_GRAY),
    ("  traffic realization gets a lower average score", 11, False, LIGHT_GRAY),
    ("  than one that generalizes across patterns.", 11, False, LIGHT_GRAY),
    ("", 7),
    ("Score = mean(% improvements across seed42 + seed456", 11, True, ACCENT_TEAL),
    ("         × 2 workloads = 4 simulation runs)", 11, False, ACCENT_TEAL),
])

for i, (seed, color, lines) in enumerate([
    ("Seed 42", ACCENT_TEAL, [
        ("Standard traffic realization", 12, True, WHITE),
        ("  •  Normal burstiness patterns", 11, False, LIGHT_GRAY),
        ("  •  Balanced prefix group arrivals", 11, False, LIGHT_GRAY),
        ("  •  Representative of typical load", 11, False, LIGHT_GRAY),
        ("", 6),
        ("glia_40qps:    4314 ms baseline E2E", 11, False, MID_GRAY),
        ("prefix_heavy:   790 ms baseline E2E", 11, False, MID_GRAY),
    ]),
    ("Seed 456", ACCENT_BLUE, [
        ("Stress-test traffic realization", 12, True, WHITE),
        ("  •  Different burst timing patterns", 11, False, LIGHT_GRAY),
        ("  •  Slightly different prefix skew", 11, False, LIGHT_GRAY),
        ("  •  Tests robustness of the algorithm", 11, False, LIGHT_GRAY),
        ("", 6),
        ("Ensures generalization, not lucky draws", 11, False, MID_GRAY),
        ("Robustness also validated w/ seed 789", 11, False, MID_GRAY),
    ]),
]):
    sy = Inches(1.42) + i * Inches(2.65)
    R(sl, Inches(6.8), sy, Inches(5.95), Inches(2.48), fill=CARD_BG, line=color)
    BAR(sl, Inches(6.8), sy, Inches(5.95), Inches(0.055), color)
    T(sl, Inches(7.0), sy+Inches(0.1), Inches(5.6), Inches(0.36),
      seed, size=16, bold=True, color=color)
    MT(sl, Inches(7.0), sy+Inches(0.52), Inches(5.6), Inches(1.85), lines)

ftr(sl, "4 simulation runs per iteration: seed42×glia_40qps, seed42×prefix_heavy, seed456×glia_40qps, seed456×prefix_heavy")

# ════════════════════════════════════════════════════════
# SLIDE 4B — Seeds (simple diagram)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); vtag(sl, "B")
hdr(sl, "Why Two Seeds?",
    "To make sure the algorithm works in general — not just for one lucky traffic draw")

R(sl, Inches(0.6), Inches(1.42), Inches(12.1), Inches(0.85), fill=CARD_BG)
T(sl, Inches(0.8), Inches(1.52), Inches(11.5), Inches(0.65),
  "Analogy: You wouldn't judge a doctor on a single patient.\n"
  "We test every routing algorithm on two different random traffic scenarios, then average.",
  size=13, color=LIGHT_GRAY)

# Flow diagram
dy = Inches(2.5)
R(sl, Inches(0.5), dy, Inches(2.5), Inches(1.1), fill=CARD_BG, line=ACCENT_TEAL)
T(sl, Inches(0.65), dy+Inches(0.18), Inches(2.2), Inches(0.75),
  "Evolved\nAlgorithm", size=14, bold=True, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)
T(sl, Inches(3.1), dy+Inches(0.0),  Inches(0.45), Inches(0.4), "→", size=18, color=MID_GRAY)
T(sl, Inches(3.1), dy+Inches(0.65), Inches(0.45), Inches(0.4), "→", size=18, color=MID_GRAY)
R(sl, Inches(3.65), dy-Inches(0.3), Inches(3.55), Inches(1.05), fill=CARD_BG, line=ACCENT_ORANGE)
T(sl, Inches(3.82), dy-Inches(0.22), Inches(3.22), Inches(0.9),
  "Seed 42: glia_40qps\n          prefix_heavy", size=12, color=ACCENT_ORANGE)
R(sl, Inches(3.65), dy+Inches(0.9), Inches(3.55), Inches(1.05), fill=CARD_BG, line=ACCENT_BLUE)
T(sl, Inches(3.82), dy+Inches(0.98), Inches(3.22), Inches(0.9),
  "Seed 456: glia_40qps\n           prefix_heavy", size=12, color=ACCENT_BLUE)
T(sl, Inches(7.32), dy+Inches(0.0),  Inches(0.45), Inches(0.4), "→", size=18, color=MID_GRAY)
T(sl, Inches(7.32), dy+Inches(0.65), Inches(0.45), Inches(0.4), "→", size=18, color=MID_GRAY)
R(sl, Inches(7.9), dy, Inches(2.7), Inches(1.1), fill=CARD_BG, line=GREEN)
T(sl, Inches(8.05), dy+Inches(0.08), Inches(2.42), Inches(0.95),
  "Average score\n(4 runs total)", size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
T(sl, Inches(10.72), dy+Inches(0.35), Inches(0.45), Inches(0.4), "→", size=18, color=MID_GRAY)
R(sl, Inches(11.25), dy, Inches(1.6), Inches(1.1), fill=CARD_BG, line=ACCENT_TEAL)
T(sl, Inches(11.35), dy+Inches(0.08), Inches(1.42), Inches(0.95),
  "+11.5%\nfinal score", size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

MT(sl, Inches(0.6), Inches(3.9), Inches(12.1), Inches(3.1), [
    ("Why does this matter?", 15, True, WHITE),
    ("", 7),
    ("Without multi-seed, an algorithm could 'luck out' on a single traffic pattern "
     "and score well without actually being better.", 13, False, LIGHT_GRAY),
    ("", 6),
    ("With two seeds, the algorithm must prove it works across different traffic realizations.", 13, False, LIGHT_GRAY),
    ("", 6),
    ("In this experiment: seeds 42 and 456 both show consistent improvement — the +11.5% is real.", 13, True, GREEN),
    ("", 6),
    ("Note: post-experiment robustness validation also tested seed 789 — results remained consistent.", 11, False, MID_GRAY),
], ds=13)
ftr(sl, "Each score = mean of % improvements across 4 runs  •  Robustness validated with 3 seeds post-experiment")


# ════════════════════════════════════════════════════════
# SLIDE 5A — Mutator / EVOLVE-BLOCK (code diff)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); vtag(sl, "A")
hdr(sl, "The Mutator: What the LLM Actually Changed",
    "Only the EVOLVE-BLOCK region is editable — the rest of the router stays fixed")

T(sl, Inches(0.6), Inches(1.42), Inches(5.85), Inches(0.32),
  "BEFORE  (initial program — simple fixed weights)", size=12, bold=True, color=MID_GRAY)
R(sl, Inches(0.6), Inches(1.8), Inches(5.85), Inches(4.72),
  fill=RGBColor(0x12, 0x12, 0x22))
MT(sl, Inches(0.78), Inches(1.9), Inches(5.55), Inches(4.5), [
    ("// EVOLVE-BLOCK-START", 9, False, MID_GRAY),
    ("scores := make(map[string]float64, ...)", 9, False, LIGHT_GRAY),
    ("for i, scorer := range ws.scorers {", 9, False, LIGHT_GRAY),
    ("    dimScores := scorer(req, snapshots)", 9, False, LIGHT_GRAY),
    ("    for _, snap := range snapshots {", 9, False, LIGHT_GRAY),
    ("        s := dimScores[snap.ID]", 9, False, LIGHT_GRAY),
    ("        if s < 0 { s = 0 }  // clamp", 9, False, LIGHT_GRAY),
    ("        if s > 1 { s = 1 }", 9, False, LIGHT_GRAY),
    ("        scores[snap.ID] += s * ws.weights[i]", 9, False, LIGHT_GRAY),
    ("    }", 9, False, LIGHT_GRAY),
    ("}", 9, False, LIGHT_GRAY),
    ("", 5),
    ("// simple argmax", 9, False, LIGHT_GRAY),
    ("bestScore := -1.0", 9, False, LIGHT_GRAY),
    ("for _, snap := range snapshots {", 9, False, LIGHT_GRAY),
    ("    if scores[snap.ID] > bestScore {", 9, False, LIGHT_GRAY),
    ("        bestScore = scores[snap.ID]", 9, False, LIGHT_GRAY),
    ("    }", 9, False, LIGHT_GRAY),
    ("}", 9, False, LIGHT_GRAY),
    ("// ... tie-breaking ...", 9, False, LIGHT_GRAY),
    ("// EVOLVE-BLOCK-END", 9, False, MID_GRAY),
], ds=9)

T(sl, Inches(6.82), Inches(1.42), Inches(5.85), Inches(0.32),
  "AFTER  (evolved by LLM at iteration 16)", size=12, bold=True, color=ACCENT_TEAL)
R(sl, Inches(6.82), Inches(1.8), Inches(5.85), Inches(4.72),
  fill=RGBColor(0x12, 0x12, 0x22))
MT(sl, Inches(7.0), Inches(1.9), Inches(5.55), Inches(4.5), [
    ("// EVOLVE-BLOCK-START", 9, False, MID_GRAY),
    ("// Pre-compute all scorer outputs", 9, False, ACCENT_TEAL),
    ("allDimScores := make([]map[string]float64,...)", 9, False, GREEN),
    ("for i, scorer := range ws.scorers { ... }", 9, False, GREEN),
    ("", 5),
    ("// Fresh load signal — best for adaptive weighting", 9, False, ACCENT_TEAL),
    ("minInflight := ... // min across all instances", 9, False, GREEN),
    ("", 5),
    ("// Find best prefix-affinity instance", 9, False, ACCENT_TEAL),
    ("bestPrefixID, bestPrefixScore := ...", 9, False, GREEN),
    ("", 5),
    ("// Adaptive weights: decay prefix when overloaded", 9, False, ACCENT_TEAL),
    ("aw := copy(ws.weights)", 9, False, GREEN),
    ("if delta := cachedLoad - minInflight; delta > 0 {", 9, False, GREEN),
    ("  decay := 1.0 / (1.0 + 0.6*float64(delta))", 9, False, GREEN),
    ("  aw[0] = ws.weights[0] * decay // reduce prefix", 9, False, GREEN),
    ("  aw[1] = 1.0 - aw[0]          // boost load-bal", 9, False, GREEN),
    ("}", 9, False, GREEN),
    ("// KV pressure penalty (>90% util)", 9, False, ACCENT_TEAL),
    ("scores[id] -= 0.5*(KVUtil-0.9)/0.1  // if >90%", 9, False, GREEN),
    ("// Fresh tiebreaker: +0.01/(1+InFlight)", 9, False, GREEN),
    ("// EVOLVE-BLOCK-END", 9, False, MID_GRAY),
], ds=9)
ftr(sl, "The EVOLVE-BLOCK marks the mutable region. LLM replaces everything between the markers. All else is locked.")

# ════════════════════════════════════════════════════════
# SLIDE 5B — Mutator (3 innovations, conceptual)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); vtag(sl, "B")
hdr(sl, "Three Things the LLM Discovered",
    "The LLM rewrote the routing logic — here's what changed and why it helps")

R(sl, Inches(0.6), Inches(1.42), Inches(12.1), Inches(0.72), fill=CARD_BG)
T(sl, Inches(0.8), Inches(1.52), Inches(11.5), Inches(0.55),
  "The routing function has an 'editable zone' (EVOLVE-BLOCK). Everything inside it can be rewritten.\n"
  "Everything outside — the API, config loading, scorer pipeline — is locked and unchanged.",
  size=12, color=LIGHT_GRAY)

iw = Inches(3.85); ih = Inches(4.3); iy = Inches(2.3)
for i, (tag, title, accent, desc) in enumerate([
    ("Innovation 1", "Adaptive Prefix-Weight Decay", ACCENT_TEAL,
     "When the instance with the best cached prefix is overloaded "
     "(high InFlightRequests), reduce its routing weight.\n\n"
     "Formula: decay = 1 / (1 + 0.6 × load_delta)\n\n"
     "This prevents piling requests onto one hot instance just "
     "because it has a useful cache. The router balances cache "
     "benefit against load fairly."),
    ("Innovation 2", "KV Memory Pressure Penalty", ACCENT_ORANGE,
     "When an instance's KV cache is more than 90% full, "
     "subtract from its score.\n\n"
     "Formula: penalty = -0.5 × (KVUtil - 0.9) / 0.1\n\n"
     "Prevents routing to instances about to hit memory limits, "
     "which cause slow cache evictions and performance degradation."),
    ("Innovation 3", "Fresh Load Tiebreaker", ACCENT_PURPLE,
     "Use InFlightRequests (updated every call, zero lag) to "
     "break ties between equally-scored instances.\n\n"
     "Formula: bonus = +0.01 / (1 + InFlightRequests)\n\n"
     "Other load signals (QueueDepth, BatchSize) are 5 seconds "
     "stale. InFlightRequests is always fresh — it is the "
     "best tiebreaker available."),
]):
    ix = Inches(0.5) + i * (iw + Inches(0.27))
    R(sl, ix, iy, iw, ih, fill=CARD_BG, line=accent)
    BAR(sl, ix, iy, iw, Inches(0.055), accent)
    T(sl, ix+Inches(0.15), iy+Inches(0.1), iw-Inches(0.3), Inches(0.26),
      tag, size=11, bold=True, color=accent)
    T(sl, ix+Inches(0.15), iy+Inches(0.4), iw-Inches(0.3), Inches(0.38),
      title, size=13, bold=True, color=WHITE)
    T(sl, ix+Inches(0.15), iy+Inches(0.85), iw-Inches(0.3), ih-Inches(1.0),
      desc, size=11, color=LIGHT_GRAY)
ftr(sl, "All three independently discovered by the LLM. The human oracle (v14) also uses adaptive weighting — different math, same insight.")


# ════════════════════════════════════════════════════════
# SLIDE 6A — Prompt (full detail)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); vtag(sl, "A")
hdr(sl, "Prompt Summary: What We Told the LLM",
    "System message in config.yaml — context the LLM receives before writing code")

R(sl, Inches(0.6), Inches(1.42), Inches(6.15), Inches(5.6), fill=CARD_BG, line=ACCENT_BLUE)
T(sl, Inches(0.8), Inches(1.55), Inches(5.9), Inches(0.35),
  "System Prompt Sections", size=14, bold=True, color=ACCENT_BLUE)
MT(sl, Inches(0.8), Inches(2.0), Inches(5.9), Inches(4.85), [
    ("ROLE", 12, True, WHITE),
    ("  Optimize a routing function for a 4-instance LLM cluster.", 11, False, LIGHT_GRAY),
    ("  Goal: minimize E2E latency across two workloads.", 11, False, LIGHT_GRAY),
    ("", 7),
    ("BASELINE", 12, True, WHITE),
    ("  1:1 router — fixed 0.5 prefix-affinity + 0.5 load-balance.", 11, False, LIGHT_GRAY),
    ("  Strong on general traffic; vulnerable on skewed prefix.", 11, False, LIGHT_GRAY),
    ("", 7),
    ("SCORERS", 12, True, WHITE),
    ("  ws.scorers[0] — prefix-affinity: ~1.0 for cached instance. FRESH.", 11, False, LIGHT_GRAY),
    ("  ws.scorers[1] — load-balance: 1/(1+EffectiveLoad()). STALE.", 11, False, LIGHT_GRAY),
    ("", 7),
    ("KEY SIGNALS", 12, True, WHITE),
    ("  snap.InFlightRequests — FRESH. Best load signal. Use this.", 11, True, GREEN),
    ("  snap.KVUtilization — 5s stale. Good safety filter.", 11, False, ACCENT_ORANGE),
    ("  snap.QueueDepth, BatchSize — 5s stale. Avoid as primary.", 11, False, RED),
    ("", 7),
    ("OPPORTUNITY", 12, True, WHITE),
    ("  prefix_heavy has 6 prefix groups, skewed (45/18/12/10/8/7%).", 11, False, LIGHT_GRAY),
    ("  Fixed weights create hotspot. Adaptive algo can win by 15%+.", 11, False, LIGHT_GRAY),
    ("", 7),
    ("RULES", 12, True, WHITE),
    ("  1. Only modify EVOLVE-BLOCK  2. Valid Go  3. All instances scored", 11, False, LIGHT_GRAY),
    ("  4. Don't remove load-balancing (prevents starvation)", 11, False, LIGHT_GRAY),
])

R(sl, Inches(7.1), Inches(1.42), Inches(5.9), Inches(5.6), fill=CARD_BG, line=ACCENT_TEAL)
T(sl, Inches(7.3), Inches(1.55), Inches(5.6), Inches(0.35),
  "Design Choices & Rationale", size=14, bold=True, color=ACCENT_TEAL)
MT(sl, Inches(7.3), Inches(2.0), Inches(5.6), Inches(4.85), [
    ("Why mention the OPPORTUNITY?", 12, True, WHITE),
    ("  Giving a concrete hint (hotspot problem) dramatically", 11, False, LIGHT_GRAY),
    ("  speeds convergence. Without it, LLM wastes iterations", 11, False, LIGHT_GRAY),
    ("  on random mutations that miss the key insight.", 11, False, LIGHT_GRAY),
    ("", 7),
    ("Why label FRESH vs STALE signals?", 12, True, WHITE),
    ("  Stale signals are misleading under bursty load.", 11, False, LIGHT_GRAY),
    ("  Telling the LLM which signals to trust directly led", 11, False, LIGHT_GRAY),
    ("  to the InFlightRequests tiebreaker innovation.", 11, False, LIGHT_GRAY),
    ("", 7),
    ("Why the no-starvation rule?", 12, True, WHITE),
    ("  Without it, LLMs over-specialize — ignoring load-balance.", 11, False, LIGHT_GRAY),
    ("  This rule caused 100% build success (no degenerate code).", 11, False, LIGHT_GRAY),
    ("", 7),
    ("Score formula shown to LLM:", 12, True, WHITE),
    ("  score = mean % improvement over baseline", 11, True, ACCENT_TEAL),
    ("  across both workloads. Positive = better.", 11, False, MID_GRAY),
    ("", 7),
    ("Temperature: 1.0 throughout (encourages diversity)", 11, False, MID_GRAY),
    ("Max tokens: 32,000 (full Go code fits)", 11, False, MID_GRAY),
])
ftr(sl, "Full system prompt in benchmarks/blis_router/config.yaml  •  OpenEvolve also provides top-k examples from its archive")

# ════════════════════════════════════════════════════════
# SLIDE 6B — Prompt (simplified 5-card layout)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); vtag(sl, "B")
hdr(sl, "What We Told the LLM",
    "Simple English version of the system prompt — 5 key pieces of context")

prompt_cards = [
    ("Your job", ACCENT_BLUE,
     "Improve a Go routing function for a 4-server LLM cluster. Goal: minimize request latency."),
    ("The current\nbaseline", ACCENT_TEAL,
     "Current router: 50% prefix-affinity + 50% load-balance. Works okay generally, "
     "but creates a hotspot when one prefix group dominates traffic."),
    ("Signals you\ncan use", ACCENT_ORANGE,
     "InFlightRequests — instant, always fresh. Best load signal.\n"
     "KVUtilization — 5 seconds old. Good for detecting memory pressure.\n"
     "QueueDepth / BatchSize — also 5s stale. Use only as secondary."),
    ("The\nopportunity", GREEN,
     "On the prefix-heavy workload, 45% of requests want the same server.\n"
     "This causes a hotspot. An adaptive algorithm that backs off when the "
     "cached instance is busy can win by 15% or more."),
    ("The rules", MID_GRAY,
     "1. Only change the EVOLVE-BLOCK region.\n"
     "2. Must be valid Go code (must compile).\n"
     "3. Don't eliminate load-balancing entirely."),
]
pw = Inches(4.82); ph = Inches(2.2)
for i, (title, accent, body) in enumerate(prompt_cards):
    row = i // 3; col = i % 3
    if i == 3: col = 0
    if i == 4: col = 1
    px = Inches(0.45) + col * (pw + Inches(0.22))
    py = Inches(1.42) + row * (ph + Inches(0.22))
    R(sl, px, py, pw, ph, fill=CARD_BG, line=accent)
    BAR(sl, px, py, pw, Inches(0.055), accent)
    T(sl, px+Inches(0.15), py+Inches(0.1), pw-Inches(0.3), Inches(0.4),
      title, size=13, bold=True, color=accent)
    T(sl, px+Inches(0.15), py+Inches(0.55), pw-Inches(0.3), ph-Inches(0.65),
      body, size=11, color=LIGHT_GRAY)
ftr(sl, "Designed to give just enough context — without over-constraining the search or prescribing the algorithm")


# ════════════════════════════════════════════════════════
# SLIDE 7A — Workloads (technical)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); vtag(sl, "A")
hdr(sl, "Two Evaluation Workloads",
    "glia_40qps (general traffic) and glia_prefix_heavy (skewed prefix traffic)")

ww = Inches(5.9); wh = Inches(5.4); wy = Inches(1.42)
# Workload 1
R(sl, Inches(0.6), wy, ww, wh, fill=CARD_BG, line=ACCENT_BLUE)
BAR(sl, Inches(0.6), wy, ww, Inches(0.06), ACCENT_BLUE)
T(sl, Inches(0.82), wy+Inches(0.1), ww-Inches(0.3), Inches(0.42),
  "Workload 1: glia_40qps", size=17, bold=True, color=ACCENT_BLUE)
T(sl, Inches(0.82), wy+Inches(0.58), ww-Inches(0.3), Inches(0.28),
  "General / mixed traffic", size=12, color=MID_GRAY)
MT(sl, Inches(0.82), wy+Inches(0.95), ww-Inches(0.3), Inches(4.2), [
    ("Traffic profile:", 12, True, WHITE),
    ("  •  1000 requests at 40 QPS aggregate", 11, False, LIGHT_GRAY),
    ("  •  Bursty arrivals (gamma distribution, CV=7.3)", 11, False, LIGHT_GRAY),
    ("  •  90% ShareGPT-like: ~500 prompt, ~250 decode tokens", 11, False, LIGHT_GRAY),
    ("  •  5% heavy-prompt (10× input tokens)", 11, False, LIGHT_GRAY),
    ("  •  5% heavy-decode (10× output tokens)", 11, False, LIGHT_GRAY),
    ("  •  No prefix caching — each request is unique", 11, False, LIGHT_GRAY),
    ("", 7),
    ("What it tests:", 12, True, WHITE),
    ("  Load-balance under bursty general traffic.", 11, False, LIGHT_GRAY),
    ("  No prefix structure → routing relies on load signals.", 11, False, LIGHT_GRAY),
    ("  Already near-optimal with fixed 1:1 weights.", 11, False, LIGHT_GRAY),
    ("", 7),
    ("Baseline (1:1) E2E:  4314 ms", 12, True, ACCENT_BLUE),
    ("Evolved E2E:          4303 ms  (+0.3%)", 12, False, GREEN),
    ("Challenging to improve — near physical limit", 11, False, MID_GRAY),
])

# Workload 2
R(sl, Inches(6.9), wy, ww, wh, fill=CARD_BG, line=ACCENT_ORANGE)
BAR(sl, Inches(6.9), wy, ww, Inches(0.06), ACCENT_ORANGE)
T(sl, Inches(7.12), wy+Inches(0.1), ww-Inches(0.3), Inches(0.42),
  "Workload 2: glia_prefix_heavy", size=17, bold=True, color=ACCENT_ORANGE)
T(sl, Inches(7.12), wy+Inches(0.58), ww-Inches(0.3), Inches(0.28),
  "Prefix-heavy / skewed traffic", size=12, color=MID_GRAY)
MT(sl, Inches(7.12), wy+Inches(0.95), ww-Inches(0.3), Inches(4.2), [
    ("Traffic profile:", 12, True, WHITE),
    ("  •  1500 requests at 85 QPS aggregate", 11, False, LIGHT_GRAY),
    ("  •  6 prefix groups (A–F), 14,336-token shared prefixes", 11, False, LIGHT_GRAY),
    ("  •  Group A dominant: 45% of all requests", 11, True, ACCENT_ORANGE),
    ("  •  Groups B-E: 8–18% each", 11, False, LIGHT_GRAY),
    ("  •  Group F: 7%, no prefix (random traffic)", 11, False, LIGHT_GRAY),
    ("  •  Mixed SLO classes (batch / interactive / realtime)", 11, False, LIGHT_GRAY),
    ("", 7),
    ("What it tests:", 12, True, WHITE),
    ("  Prefix-affinity routing under heavy traffic skew.", 11, False, LIGHT_GRAY),
    ("  Group A (45%) creates hotspot on instance 1.", 11, False, LIGHT_GRAY),
    ("  Fixed weights fail here — adaptive weights win.", 11, False, LIGHT_GRAY),
    ("", 7),
    ("Baseline (1:1) E2E:  790 ms", 12, True, ACCENT_ORANGE),
    ("Evolved E2E:          700 ms  (+11.4%)", 12, True, GREEN),
    ("← This is where the big improvement comes from!", 11, True, ACCENT_ORANGE),
])
ftr(sl, "Both workloads run in BLIS discrete-event simulator  •  Model: qwen2.5-7b-instruct on 4×H100 GPUs  •  vLLM v0.11.0")

# ════════════════════════════════════════════════════════
# SLIDE 7B — Workloads (visual story)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); vtag(sl, "B")
hdr(sl, "Two Workloads — Two Different Stories",
    "One tests general balance. The other tests hotspot handling under skewed traffic.")

R(sl, Inches(0.6), Inches(1.42), Inches(12.1), Inches(0.72), fill=CARD_BG)
T(sl, Inches(0.8), Inches(1.52), Inches(11.5), Inches(0.55),
  "Think of 4 servers as 4 cafeterias.  "
  "Workload 1 = customers choosing randomly.  "
  "Workload 2 = 45% of all customers want the same cafeteria.",
  size=13, color=LIGHT_GRAY)

cw2 = Inches(5.9); ch2 = Inches(4.52); cy2 = Inches(2.3)
# WL1
R(sl, Inches(0.6), cy2, cw2, ch2, fill=CARD_BG, line=ACCENT_BLUE)
T(sl, Inches(0.82), cy2+Inches(0.1), cw2-Inches(0.3), Inches(0.4),
  "glia_40qps — General Traffic", size=15, bold=True, color=ACCENT_BLUE)
# Server boxes with equal arrows
for j, sx in enumerate([Inches(0.88), Inches(1.95), Inches(3.02), Inches(4.09)]):
    R(sl, sx, cy2+Inches(0.68), Inches(0.88), Inches(0.95), fill=SUBTLE_GRAY, line=ACCENT_BLUE)
    T(sl, sx, cy2+Inches(0.68), Inches(0.88), Inches(0.95),
      "GPU\nSrv", size=9, color=WHITE, align=PP_ALIGN.CENTER)
    T(sl, sx, cy2+Inches(0.48), Inches(0.88), Inches(0.25),
      "↓ 25%", size=9, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
T(sl, Inches(0.82), cy2+Inches(1.75), cw2-Inches(0.3), Inches(2.6),
  "Requests arrive randomly. No prefix groups.\n"
  "The router just needs to spread load evenly.\n\n"
  "Already handled well by the 1:1 baseline.\n"
  "→ Only +0.3% improvement from evolved algorithm.\n\n"
  "1000 requests  •  40 QPS  •  Bursty (gamma, CV=7.3)",
  size=12, color=LIGHT_GRAY)

# WL2
R(sl, Inches(6.9), cy2, cw2, ch2, fill=CARD_BG, line=ACCENT_ORANGE)
T(sl, Inches(7.12), cy2+Inches(0.1), cw2-Inches(0.3), Inches(0.4),
  "glia_prefix_heavy — Skewed Traffic", size=15, bold=True, color=ACCENT_ORANGE)
# Server boxes with skewed arrows
for j, (sx, pct, acc) in enumerate([
    (Inches(7.18), "45%↓↓↓↓↓", ACCENT_ORANGE),
    (Inches(8.25), "18%↓↓", LIGHT_GRAY),
    (Inches(9.32), "12%↓", MID_GRAY),
    (Inches(10.39), "10%↓", MID_GRAY),
]):
    R(sl, sx, cy2+Inches(0.68), Inches(0.88), Inches(0.95), fill=SUBTLE_GRAY, line=acc)
    T(sl, sx, cy2+Inches(0.68), Inches(0.88), Inches(0.95),
      "GPU\nSrv", size=9, color=WHITE, align=PP_ALIGN.CENTER)
    T(sl, sx, cy2+Inches(0.42), Inches(0.88), Inches(0.3),
      pct, size=8, color=acc, align=PP_ALIGN.CENTER)
T(sl, Inches(7.12), cy2+Inches(1.75), cw2-Inches(0.3), Inches(2.6),
  "45% of requests all go to server 1 (prefix group A).\n"
  "This creates a hotspot — queue builds up → slow!\n\n"
  "The evolved algorithm detects the overload and\n"
  "dynamically shifts some requests to other servers.\n"
  "→ +11.4% improvement  (7× bigger win!)\n\n"
  "1500 requests  •  85 QPS  •  6 prefix groups",
  size=12, color=LIGHT_GRAY)
ftr(sl, "The prefix-heavy workload is where adaptive routing matters most — and where SkyDiscover's discovery shines")


# ════════════════════════════════════════════════════════
# SLIDE 8A — Baselines (full detail cards)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); vtag(sl, "A")
hdr(sl, "Baselines: Six Routing Algorithms We Compared Against",
    "From naive queue-depth routing to carefully hand-crafted adaptive oracle")

baselines_a = [
    ("LLQ", "Least Loaded Queue", ACCENT_ORANGE, RED, "-61% vs 1:1",
     "Routes to instance with shortest QueueDepth. Stale (5s) signal — "
     "outdated by the time routing happens. Fails badly under bursty load."),
    ("LOR", "Least Outstanding Requests", MID_GRAY, RED, "-11.6% vs 1:1",
     "Routes to instance with fewest InFlightRequests. Fresh signal, good load "
     "balance. But ignores prefix affinity entirely — poor cache hit rates."),
    ("Glia HRA", "Head-Room Allocator", ACCENT_PURPLE, RED, "-14.3% vs 1:1",
     "KV-aware: estimates free blocks, projects request cost, applies safety "
     "fraction. Good at avoiding OOM but complex calc — slightly worse than 1:1."),
    ("1:1", "Fixed 50/50 Prefix + Load-Balance", ACCENT_BLUE, WHITE, "0% (control)",
     "Our control baseline. 50% prefix-affinity + 50% load-balance, fixed weights. "
     "Simple, robust. Strong on general traffic; weak on skewed prefix workloads."),
    ("3:2:2", "Weighted 3:2:2 Prefix+LB+KV", ACCENT_TEAL, MID_GRAY, "-0.7% vs 1:1",
     "Hand-tuned weights: 3 prefix + 2 load-balance + 2 KV. Nearly identical to "
     "1:1. Weight tuning alone is insufficient — algorithm structure matters more."),
    ("Oracle v14", "Hand-crafted Adaptive Router", GREEN, GREEN, "Best human (Evolved beats by +0.6%)",
     "Best human-engineered baseline. Uses adaptive prefix decay with a hard "
     "threshold at 1.8× avg load. Very close to evolved — AI found a better math."),
]
bw = Inches(5.82); bh = Inches(0.95)
for i, (name, full, color, sc, score, desc) in enumerate(baselines_a):
    row = i // 2; col = i % 2
    bx = Inches(0.55) + col * (bw + Inches(0.6))
    by = Inches(1.42) + row * (bh + Inches(0.14))
    R(sl, bx, by, bw, bh, fill=CARD_BG)
    BAR(sl, bx, by, Inches(0.08), bh, color)
    T(sl, bx+Inches(0.2), by+Inches(0.08), Inches(1.0), Inches(0.3),
      name, size=13, bold=True, color=color)
    T(sl, bx+Inches(0.2), by+Inches(0.44), Inches(2.2), Inches(0.24),
      full, size=9, color=MID_GRAY)
    T(sl, bx+Inches(3.2), by+Inches(0.08), Inches(2.5), Inches(0.3),
      score, size=10, bold=True, color=sc, align=PP_ALIGN.RIGHT)
    T(sl, bx+Inches(0.2), by+Inches(0.68), bw-Inches(0.35), Inches(0.24),
      desc, size=9, color=LIGHT_GRAY)

R(sl, Inches(0.55), Inches(6.5), Inches(12.1), Inches(0.52), fill=CARD_BG)
T(sl, Inches(0.75), Inches(6.56), Inches(11.5), Inches(0.38),
  "Ranking (combined, lower E2E = better):   Evolved  ≈  Oracle v14  >  1:1  ≈  3:2:2  >  LOR  >  Glia  >>  LLQ",
  size=12, color=LIGHT_GRAY)
ftr(sl, "All baselines measured: same workloads, same seeds (42+456), same model (qwen_7b) as the evolved algorithm")

# ════════════════════════════════════════════════════════
# SLIDE 8B — Baselines (visual bar chart)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); vtag(sl, "B")
hdr(sl, "Baselines: The Competition Lineup",
    "Lower latency = better  •  Two workloads side by side")

bdata = [
    ("LLQ",        6357, 1300, RED),
    ("Glia",       4457, 880,  ACCENT_PURPLE),
    ("1:1 (ctrl)", 4314, 790,  ACCENT_BLUE),
    ("3:2:2",      4311, 818,  ACCENT_TEAL),
    ("Oracle v14", 4303, 706,  GREEN),
    ("Evolved ★",  4303, 700,  WHITE),
]
col_w_b = Inches(13.333) / (len(bdata) + 0.5)
# Name labels
for j, (name, e40, eph, color) in enumerate(bdata):
    bx = Inches(0.4) + j * col_w_b
    bw_b = col_w_b - Inches(0.12)
    R(sl, bx, Inches(1.42), bw_b, Inches(0.72), fill=CARD_BG,
      line=color if color != WHITE else ACCENT_TEAL)
    T(sl, bx+Inches(0.04), Inches(1.47), bw_b-Inches(0.08), Inches(0.62),
      name, size=10, bold=True,
      color=color if color != WHITE else ACCENT_TEAL,
      align=PP_ALIGN.CENTER)

# glia_40qps bars
T(sl, Inches(0.4), Inches(2.25), Inches(7), Inches(0.3),
  "glia_40qps E2E latency (ms) — lower is better", size=12, bold=True, color=WHITE)
max40 = 6500.0; bar_bot_40 = Inches(5.0); bar_max_h = Inches(2.5)
for j, (name, e40, eph, color) in enumerate(bdata):
    bx = Inches(0.4) + j * col_w_b + Inches(0.1)
    bw_b = col_w_b - Inches(0.3)
    fc = color if color != WHITE else ACCENT_TEAL
    bh_v = (e40 / max40) * bar_max_h
    by_v = bar_bot_40 - bh_v
    BAR(sl, bx, by_v, bw_b, bh_v, fc)
    T(sl, bx-Inches(0.05), by_v-Inches(0.28), bw_b+Inches(0.1), Inches(0.26),
      f"{int(e40)}", size=9, color=fc, align=PP_ALIGN.CENTER)
# reference line at 4303
ref_y = bar_bot_40 - (4303/max40)*bar_max_h
rl = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), ref_y, Inches(13.0)*0.92, Inches(0.018))
rl.fill.solid(); rl.fill.fore_color.rgb = GREEN; rl.line.fill.background()

# prefix_heavy bars
T(sl, Inches(0.4), Inches(5.12), Inches(7), Inches(0.3),
  "prefix_heavy E2E latency (ms) — lower is better", size=12, bold=True, color=WHITE)
maxph = 1400.0; bar_bot_ph = Inches(7.05); bar_max_ph = Inches(1.65)
for j, (name, e40, eph, color) in enumerate(bdata):
    bx = Inches(0.4) + j * col_w_b + Inches(0.1)
    bw_b = col_w_b - Inches(0.3)
    fc = color if color != WHITE else ACCENT_TEAL
    bh_v = (eph / maxph) * bar_max_ph
    by_v = bar_bot_ph - bh_v
    BAR(sl, bx, by_v, bw_b, bh_v, fc)
    T(sl, bx-Inches(0.05), by_v-Inches(0.28), bw_b+Inches(0.1), Inches(0.26),
      f"{int(eph)}", size=9, color=fc, align=PP_ALIGN.CENTER)
ref_y2 = bar_bot_ph - (700/maxph)*bar_max_ph
rl2 = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), ref_y2, Inches(13.0)*0.92, Inches(0.018))
rl2.fill.solid(); rl2.fill.fore_color.rgb = ACCENT_TEAL; rl2.line.fill.background()

ftr(sl, "Key insight: on glia_40qps all hand-tuned algorithms are nearly tied. Differentiation is on prefix_heavy.")


# ════════════════════════════════════════════════════════
# SLIDE 9A — Results (full table)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); vtag(sl, "A")
hdr(sl, "Results: +11.5% vs Best Hand-Tuned Baseline",
    "Per-workload breakdown  •  qwen_7b  •  seeds 42+456  •  evolved at iteration 16/50")

# Table: algorithm | g40 E2E | vs evolved | ph E2E | vs evolved | combined vs 1:1
rows = [
    ("LLQ",        "6357 ms", "+31.0%", "1300 ms", "+46.2%", "+42.0%", RED,          False),
    ("Glia HRA",   "4457 ms", "+3.6%",  " 880 ms", "+20.6%", "+17.8%", ACCENT_PURPLE,False),
    ("3:2:2",      "4311 ms", "+0.2%",  " 818 ms", "+14.4%", "+10.9%", ACCENT_TEAL,  False),
    ("1:1 (ctrl)", "4314 ms", "+0.3%",  " 790 ms", "+11.4%", "+11.5%", ACCENT_BLUE,  False),
    ("Oracle v14", "4303 ms",  "≈same", " 706 ms", "+0.9%",  "+0.6%",  GREEN,         False),
    ("Evolved ★",  "4303 ms", " —",     " 700 ms", " —",     "0% (ref)",GREEN,        True),
]
hdrs = ["Algorithm", "glia_40qps\nE2E", "vs Evolved", "prefix_heavy\nE2E", "vs Evolved", "Combined\nvs 1:1"]
cx = [Inches(0.5), Inches(2.25), Inches(3.85), Inches(5.25), Inches(6.95), Inches(8.55)]
cw_r = [Inches(1.65), Inches(1.5), Inches(1.3), Inches(1.6), Inches(1.5), Inches(1.75)]

hy = Inches(1.42)
R(sl, Inches(0.5), hy, Inches(10.6), Inches(0.48), fill=SUBTLE_GRAY, shape=MSO_SHAPE.RECTANGLE)
for j, (h, x, w) in enumerate(zip(hdrs, cx, cw_r)):
    T(sl, x+Inches(0.04), hy+Inches(0.04), w, Inches(0.4),
      h, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i, (algo, g40, g40p, ph, php, comb, color, hi) in enumerate(rows):
    ry = hy + Inches(0.5) + i * Inches(0.76)
    fill = RGBColor(0x2A, 0x2A, 0x44) if hi else CARD_BG
    R(sl, Inches(0.5), ry, Inches(10.6), Inches(0.72), fill=fill, shape=MSO_SHAPE.RECTANGLE)
    if hi:
        BAR(sl, Inches(0.5), ry, Inches(0.08), Inches(0.72), GREEN)
    vals = [algo, g40, g40p, ph, php, comb]
    for j, (val, x, w) in enumerate(zip(vals, cx, cw_r)):
        vc = WHITE if hi else (color if j == 0 else LIGHT_GRAY)
        if j in [2, 4] and "+" in str(val): vc = GREEN
        if j == 5 and "+" in str(val): vc = GREEN
        T(sl, x+Inches(0.04), ry+Inches(0.16), w, Inches(0.42),
          val, size=11, bold=hi, color=vc, align=PP_ALIGN.CENTER)

# P95 row
R(sl, Inches(0.5), Inches(6.1), Inches(10.6), Inches(0.45), fill=CARD_BG, shape=MSO_SHAPE.RECTANGLE)
T(sl, Inches(0.65), Inches(6.18), Inches(10.2), Inches(0.32),
  "P95:  glia_40qps — 1:1=17,241ms → Evolved=16,813ms (+2.5%)   |   "
  "prefix_heavy — 1:1=1,909ms → Evolved=1,435ms (+24.8%)",
  size=11, bold=True, color=ACCENT_TEAL)

# Convergence panel (right)
R(sl, Inches(11.25), Inches(1.42), Inches(1.85), Inches(5.15), fill=CARD_BG)
T(sl, Inches(11.35), Inches(1.52), Inches(1.65), Inches(0.32),
  "Score trajectory", size=10, bold=True, color=WHITE)
for k, (it, sc, c) in enumerate([
    (0,  "+0%",    MID_GRAY),
    (1,  "+7.2%",  ACCENT_BLUE),
    (2,  "+10.7%", ACCENT_TEAL),
    (16, "+11.5%", GREEN),
    (50, "+11.5%", GREEN),
]):
    ty = Inches(1.92) + k * Inches(0.88)
    T(sl, Inches(11.35), ty, Inches(1.65), Inches(0.28), f"Iter {it:2d}", size=9, color=MID_GRAY)
    T(sl, Inches(11.35), ty+Inches(0.28), Inches(1.65), Inches(0.36), sc, size=13, bold=True, color=c)

ftr(sl, "Score = 0.5×E2E + 0.5×P95 per workload, improvement = mean(1 - evolved/baseline)  •  Best at iter 16 of 50")

# ════════════════════════════════════════════════════════
# SLIDE 9B — Results (visual bars)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); vtag(sl, "B")
hdr(sl, "Results: Per-Workload Visual Comparison",
    "Lower latency = better  •  Evolved (★) vs all baselines")

T(sl, Inches(0.6), Inches(1.35), Inches(12.6), Inches(0.34),
  "Evolved:   glia_40qps = 4303 ms  (+0.3% vs 1:1)   |   "
  "prefix_heavy = 700 ms  (+11.4% vs 1:1)   |   Combined = +11.5%",
  size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

bdata9 = [
    ("LLQ",        6357, 1300, RED),
    ("Glia",       4457, 880,  ACCENT_PURPLE),
    ("1:1\n(ctrl)",4314, 790,  ACCENT_BLUE),
    ("3:2:2",      4311, 818,  ACCENT_TEAL),
    ("Oracle\nv14",4303, 706,  GREEN),
    ("Evolved\n★", 4303, 700,  WHITE),
]
barea_w = Inches(5.85); barea_h = Inches(2.7)

for chart_idx, (title, vals_key, max_v, ref_v, ref_c, bx_start) in enumerate([
    ("glia_40qps — E2E Latency (ms)", 0, 6500, 4303, GREEN, Inches(0.6)),
    ("prefix_heavy — E2E Latency (ms)", 1, 1400, 700, ACCENT_TEAL, Inches(6.9)),
]):
    T(sl, bx_start, Inches(1.82), barea_w, Inches(0.32),
      title, size=13, bold=True, color=WHITE)
    bar_bot = Inches(4.9)
    bar_slot = barea_w / len(bdata9)
    for j, (name, e40, eph, color) in enumerate(bdata9):
        val = e40 if vals_key == 0 else eph
        fc = color if color != WHITE else ACCENT_TEAL
        bh_v = (val / max_v) * barea_h
        by_v = bar_bot - bh_v
        bx_j = bx_start + j * bar_slot + Inches(0.06)
        bw_j = bar_slot - Inches(0.14)
        BAR(sl, bx_j, by_v, bw_j, bh_v, fc)
        T(sl, bx_j-Inches(0.02), bar_bot+Inches(0.04), bw_j+Inches(0.04), Inches(0.45),
          name, size=9, color=fc, align=PP_ALIGN.CENTER)
        T(sl, bx_j-Inches(0.02), by_v-Inches(0.26), bw_j+Inches(0.04), Inches(0.24),
          f"{int(val)}", size=9, color=fc, align=PP_ALIGN.CENTER)
    # reference line
    ref_y = bar_bot - (ref_v / max_v) * barea_h
    rl = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx_start, ref_y, barea_w, Inches(0.018))
    rl.fill.solid(); rl.fill.fore_color.rgb = ref_c; rl.line.fill.background()

R(sl, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.25), fill=CARD_BG)
MT(sl, Inches(0.8), Inches(5.6), Inches(11.5), Inches(1.05), [
    ("glia_40qps: all hand-tuned algorithms are already near-optimal. "
     "The evolved algorithm matches the very best (Oracle v14).", 12, False, LIGHT_GRAY),
    ("prefix_heavy: adaptive prefix-weight decay delivers a clear win. "
     "AI independently discovered the same core insight as the human oracle — but with better math.", 12, False, LIGHT_GRAY),
], ds=12)
ftr(sl, "★ = best_program.go from iter 16  •  P95: glia_40qps +2.5%, prefix_heavy +24.8%  •  0 build errors")


# ════════════════════════════════════════════════════════
# SLIDE 10A — Sim2Real (technical detail)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); vtag(sl, "A")
hdr(sl, "Sim2Real Transfer Package",
    "sim2real/blis_router/ — everything needed to validate and deploy the evolved router")

R(sl, Inches(0.6), Inches(1.42), Inches(5.85), Inches(5.6), fill=CARD_BG, line=ACCENT_TEAL)
T(sl, Inches(0.82), Inches(1.55), Inches(5.55), Inches(0.36),
  "Package Contents", size=14, bold=True, color=ACCENT_TEAL)
MT(sl, Inches(0.82), Inches(2.0), Inches(5.55), Inches(4.85), [
    ("sim2real/blis_router/", 12, True, WHITE),
    ("  README.md — instructions & sim2real story", 11, False, LIGHT_GRAY),
    ("  llm_config.yaml — hardware & cluster config", 11, False, LIGHT_GRAY),
    ("  repro.py — standalone reproduction script", 11, False, LIGHT_GRAY),
    ("", 6),
    ("  workloads/", 12, True, ACCENT_ORANGE),
    ("    workload_glia_40qps.yaml", 11, False, LIGHT_GRAY),
    ("    workload_glia_prefix_heavy.yaml", 11, False, LIGHT_GRAY),
    ("", 6),
    ("  baselines/", 12, True, ACCENT_PURPLE),
    ("    baseline_1_1.go  — 50/50 control", 11, False, LIGHT_GRAY),
    ("    baseline_llq.go  — LLQ", 11, False, LIGHT_GRAY),
    ("    baseline_glia.go — Glia HRA", 11, False, LIGHT_GRAY),
    ("    baseline_lor.go  — LOR", 11, False, LIGHT_GRAY),
    ("", 6),
    ("  best/", 12, True, GREEN),
    ("    best_program.go        — AI router", 11, False, LIGHT_GRAY),
    ("    best_program_info.json — metrics & lineage", 11, False, LIGHT_GRAY),
    ("", 6),
    ("  routing_config/routing_policy.yaml", 11, False, LIGHT_GRAY),
    ("  others/hardware_config.json, evaluator.py", 11, False, LIGHT_GRAY),
    ("  others/baseline_comparison.json", 11, False, LIGHT_GRAY),
])

R(sl, Inches(6.75), Inches(1.42), Inches(6.05), Inches(2.55), fill=CARD_BG, line=ACCENT_ORANGE)
T(sl, Inches(6.95), Inches(1.55), Inches(5.8), Inches(0.36),
  "Critical for Sim2Real Transfer", size=14, bold=True, color=ACCENT_ORANGE)
MT(sl, Inches(6.95), Inches(2.0), Inches(5.8), Inches(1.82), [
    ("a) Workloads — must match exactly", 12, True, WHITE),
    ("   QPS, prefix sizes, token dists, SLO classes", 11, False, LIGHT_GRAY),
    ("b) LLM config — Qwen2.5-7B, 4×H100, TP=1, vLLM v0.11", 12, True, WHITE),
    ("   Different model → different absolute latencies", 11, False, LIGHT_GRAY),
    ("c) Both algorithms — baseline_1_1.go AND best_program.go", 12, True, WHITE),
    ("   Run both → compute improvement vs your control", 11, False, LIGHT_GRAY),
])

R(sl, Inches(6.75), Inches(4.12), Inches(6.05), Inches(2.9), fill=CARD_BG, line=GREEN)
T(sl, Inches(6.95), Inches(4.25), Inches(5.8), Inches(0.36),
  "Expected Outcomes", size=14, bold=True, color=GREEN)
MT(sl, Inches(6.95), Inches(4.7), Inches(5.8), Inches(2.2), [
    ("Workload         Metric      Baseline  Evolved    Gain", 10, True, WHITE),
    ("─" * 48, 9, False, MID_GRAY),
    ("glia_40qps     E2E mean   ~4314 ms  ~4303 ms   +0.3%", 11, False, LIGHT_GRAY),
    ("glia_40qps     E2E P95    ~17241ms  ~16813ms   +2.5%", 11, False, LIGHT_GRAY),
    ("prefix_heavy   E2E mean    ~790 ms   ~700 ms  +11.4%", 11, True, GREEN),
    ("prefix_heavy   E2E P95    ~1909 ms  ~1435 ms  +24.8%", 11, True, GREEN),
    ("", 5),
    ("Combined improvement (evaluator formula):  ~+11.5%", 12, True, GREEN),
])
ftr(sl, "Viability: HIGH — InFlightRequests (router-local, FRESH) + KVUtil (Prometheus ~5s) — all available in real llm-d")

# ════════════════════════════════════════════════════════
# SLIDE 10B — Sim2Real (story-driven, 3 columns)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl); vtag(sl, "B")
hdr(sl, "From Simulation to Real Deployment",
    "The sim2real package makes it easy to reproduce results and validate on a real cluster")

R(sl, Inches(0.6), Inches(1.42), Inches(12.1), Inches(0.72), fill=CARD_BG)
T(sl, Inches(0.8), Inches(1.52), Inches(11.5), Inches(0.55),
  "The evolved router was discovered in simulation. The sim2real package bundles everything "
  "you need to run both the baseline and evolved router on a real cluster — and verify the gains hold.",
  size=13, color=LIGHT_GRAY)

tc_w = Inches(3.85); tc_h = Inches(4.52); tc_y = Inches(2.32)
for i, (title, accent, lines) in enumerate([
    ("What to Run", ACCENT_TEAL, [
        ("1.  Run baseline_1_1.go (control)", 12, True, WHITE),
        ("    Standard 50/50 prefix/load router.", 11, False, LIGHT_GRAY),
        ("    Establish your control latency numbers.", 11, False, LIGHT_GRAY),
        ("", 7),
        ("2.  Run best_program.go (evolved)", 12, True, WHITE),
        ("    Drop-in replacement for routing.go.", 11, False, LIGHT_GRAY),
        ("    Same config, same workloads.", 11, False, LIGHT_GRAY),
        ("", 7),
        ("3.  Compare the two", 12, True, WHITE),
        ("    The repro.py script prints a", 11, False, LIGHT_GRAY),
        ("    comparison table automatically.", 11, False, LIGHT_GRAY),
    ]),
    ("Three Critical Things", ACCENT_ORANGE, [
        ("a)  Use the exact workload specs", 12, True, WHITE),
        ("    QPS, prefix sizes, token distributions,", 11, False, LIGHT_GRAY),
        ("    SLO classes — all matter for sim fidelity.", 11, False, LIGHT_GRAY),
        ("", 7),
        ("b)  Match the LLM & hardware", 12, True, WHITE),
        ("    Qwen2.5-7B-Instruct on 4×H100, TP=1,", 11, False, LIGHT_GRAY),
        ("    vLLM v0.11.0. Different model = different", 11, False, LIGHT_GRAY),
        ("    absolute latency — gains still hold.", 11, False, LIGHT_GRAY),
        ("", 7),
        ("c)  Always compare vs baseline_1_1.go", 12, True, WHITE),
        ("    Don't judge absolute ms — judge the delta.", 11, False, LIGHT_GRAY),
    ]),
    ("What to Expect", GREEN, [
        ("glia_40qps workload:", 12, True, WHITE),
        ("  Baseline: ~4314 ms E2E", 11, False, LIGHT_GRAY),
        ("  Evolved:  ~4303 ms E2E", 11, False, LIGHT_GRAY),
        ("  Gain:     +0.3%  (small, already optimal)", 11, False, MID_GRAY),
        ("", 7),
        ("prefix_heavy workload:", 12, True, WHITE),
        ("  Baseline: ~790 ms E2E", 11, False, LIGHT_GRAY),
        ("  Evolved:  ~700 ms E2E", 11, True, GREEN),
        ("  Gain:     +11.4%  (large!)", 12, True, GREEN),
        ("", 7),
        ("All signals are production-available:", 12, True, WHITE),
        ("  InFlightRequests — router-local, FRESH", 11, False, LIGHT_GRAY),
        ("  KVUtilization — Prometheus, ~5s stale", 11, False, LIGHT_GRAY),
        ("  Prefix-affinity — router LRU, FRESH", 11, False, LIGHT_GRAY),
    ]),
]):
    x = Inches(0.5) + i * (tc_w + Inches(0.27))
    R(sl, x, tc_y, tc_w, tc_h, fill=CARD_BG, line=accent)
    BAR(sl, x, tc_y, tc_w, Inches(0.055), accent)
    T(sl, x+Inches(0.15), tc_y+Inches(0.1), tc_w-Inches(0.3), Inches(0.36),
      title, size=14, bold=True, color=accent)
    MT(sl, x+Inches(0.15), tc_y+Inches(0.55), tc_w-Inches(0.3), Inches(3.82), lines)
ftr(sl, "Deployment viability: HIGH  •  No simulator-only constructs used  •  All signals available in real llm-d")


# ════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════
out_path = os.path.join(os.path.dirname(__file__), "blis_openevolve_experiment.pptx")
prs.save(out_path)
print(f"Saved {len(prs.slides)} slides → {out_path}")
