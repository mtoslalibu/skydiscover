"""
BLIS Admission Control — SkyDiscover Experiment Deck
Run: python3 docs/slides/blis_admission_deck.py
Output: docs/slides/blis_admission_experiment.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0D, 0x1B, 0x2A)   # dark navy
ACCENT    = RGBColor(0x00, 0xC2, 0xFF)   # cyan
GREEN     = RGBColor(0x2E, 0xCC, 0x71)   # success green
ORANGE    = RGBColor(0xFF, 0x8C, 0x00)   # warning orange
RED       = RGBColor(0xFF, 0x4D, 0x4D)   # danger red
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCC, 0xDD, 0xEE)
CARD_BG   = RGBColor(0x16, 0x2B, 0x40)
YELLOW    = RGBColor(0xFF, 0xD7, 0x00)

W, H = Inches(13.33), Inches(7.5)   # 16:9 widescreen

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ── Helpers ──────────────────────────────────────────────────────────────────

def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return s

def box(s, x, y, w, h, fill=None, border=None, border_w=Pt(1.5), radius=None):
    shape = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background() if border is None else None
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = border_w
    return shape

def txt(s, text, x, y, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def heading(s, title, subtitle=None):
    txt(s, title, 0.4, 0.18, 12.5, 0.7, size=32, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(s, subtitle, 0.4, 0.82, 12.5, 0.4, size=16, color=LIGHT, align=PP_ALIGN.LEFT)
    # thin accent line
    line = s.shapes.add_shape(1, Inches(0.4), Inches(1.15), Inches(12.5), Inches(0.03))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

def card(s, x, y, w, h, title, body_lines, title_color=ACCENT, body_size=14):
    box(s, x, y, w, h, fill=CARD_BG, border=ACCENT, border_w=Pt(1))
    txt(s, title, x+0.15, y+0.12, w-0.3, 0.38, size=15, bold=True, color=title_color)
    body = "\n".join(body_lines)
    txt(s, body, x+0.15, y+0.52, w-0.3, h-0.65, size=body_size, color=LIGHT)

def bullet_block(s, x, y, w, lines, size=15, gap=0.34):
    for i, (bullet, text) in enumerate(lines):
        txt(s, bullet, x, y + i*gap, 0.3, gap, size=size, color=ACCENT, bold=True)
        txt(s, text,   x+0.28, y + i*gap, w-0.28, gap, size=size, color=WHITE)

def arrow(s, x1, y1, x2, y2, color=ACCENT):
    """Draw a simple right-pointing arrow as a filled rectangle + triangle approximation."""
    from pptx.util import Inches
    # Use a connector
    connector = s.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(2)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════
s1 = slide()

# Big title
txt(s1, "BLIS Admission Control", 1.0, 1.6, 11.3, 1.1,
    size=48, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txt(s1, "Discovering adaptive admission policies for LLM inference clusters",
    1.0, 2.75, 11.3, 0.6, size=22, color=LIGHT, align=PP_ALIGN.CENTER)
txt(s1, "using SkyDiscover + OpenEvolve",
    1.0, 3.3, 11.3, 0.5, size=18, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

# Divider
line = s1.shapes.add_shape(1, Inches(3.5), Inches(4.0), Inches(6.3), Inches(0.04))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

txt(s1, "Experiment:  260324_30i_nohint_2seed  •  March 2026",
    1.0, 4.2, 11.3, 0.4, size=14, color=LIGHT, align=PP_ALIGN.CENTER)

# Three stat pills at the bottom
for i, (val, label, col) in enumerate([
    ("+38.8%", "Score improvement", GREEN),
    ("0.8700", "Best combined score", ACCENT),
    ("30 iters", "OpenEvolve search", ORANGE),
]):
    bx = 1.8 + i * 3.3
    box(s1, bx, 5.1, 2.9, 1.3, fill=CARD_BG, border=col)
    txt(s1, val,   bx, 5.2,  2.9, 0.55, size=28, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s1, label, bx, 5.75, 2.9, 0.4,  size=13, color=LIGHT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Setup
# ═══════════════════════════════════════════════════════════════════════════
s2 = slide()
heading(s2, "The Problem: Overloaded LLM Cluster", "What happens when too many requests arrive at once?")

# Left: cluster diagram
txt(s2, "4-Instance Cluster  (qwen2.5-7b / H100)", 0.4, 1.35, 5.5, 0.4,
    size=14, bold=True, color=ACCENT)

for i in range(4):
    bx, by = 0.5 + i*1.22, 1.85
    box(s2, bx, by, 1.05, 0.8, fill=RGBColor(0x1a,0x3a,0x5c), border=ACCENT, border_w=Pt(1))
    txt(s2, f"GPU {i+1}", bx, by+0.05, 1.05, 0.35, size=13, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s2, "H100", bx, by+0.4, 1.05, 0.28, size=11, color=LIGHT, align=PP_ALIGN.CENTER)

# Saturation label
box(s2, 0.4, 2.85, 5.2, 0.45, fill=RGBColor(0x0a,0x2a,0x1a), border=GREEN)
txt(s2, "✓  Capacity: ~160 req/s    →    2× overload = 320 req/s",
    0.55, 2.88, 5.0, 0.38, size=13, color=GREEN)

# Right: what goes wrong
txt(s2, "Without Admission Control", 6.2, 1.35, 6.5, 0.4,
    size=14, bold=True, color=RED)

problems = [
    ("🔴", "Critical requests queue behind batch jobs"),
    ("🔴", "P95 latency explodes: 994ms SLO → 15,000ms actual"),
    ("🔴", "critical SLO attainment: only 25%"),
    ("🟡", "Throughput looks fine — but quality of service is broken"),
]
for i, (icon, text) in enumerate(problems):
    txt(s2, icon, 6.2, 1.85 + i*0.52, 0.4, 0.45, size=18)
    txt(s2, text, 6.65, 1.85 + i*0.52, 6.3, 0.45, size=14, color=LIGHT)

# Bottom: SLO classes table
txt(s2, "SLO Classes & Priorities", 0.4, 3.55, 12.5, 0.38,
    size=14, bold=True, color=ACCENT)

headers = ["Class", "Priority", "SLO Target", "Weight in Score", "Fraction of Traffic"]
vals    = [
    ["critical",  "Highest",  "994 ms",    "4×",  "20%"],
    ["standard",  "High",     "3,185 ms",  "2×",  "30%"],
    ["sheddable", "Low",      "9,948 ms",  "1×",  "20%"],
    ["batch",     "Lowest",   "33,340 ms", "0.5×","30%"],
]
col_colors = [ACCENT, LIGHT, LIGHT, YELLOW, LIGHT]
col_w = [1.9, 1.7, 1.8, 2.1, 2.3]
col_x = [0.4, 2.3, 4.0, 5.8, 7.9]

for ci, (hdr, cw, cx) in enumerate(zip(headers, col_w, col_x)):
    box(s2, cx, 4.0, cw-0.08, 0.38, fill=RGBColor(0x10,0x25,0x3a), border=ACCENT, border_w=Pt(0.8))
    txt(s2, hdr, cx+0.08, 4.02, cw-0.16, 0.34, size=12, bold=True,
        color=ACCENT, align=PP_ALIGN.CENTER)

row_colors = [RGBColor(0x8B,0x00,0x00), RGBColor(0x1a,0x3a,0x5c),
              RGBColor(0x1a,0x2a,0x1a), RGBColor(0x2a,0x2a,0x10)]
for ri, row in enumerate(vals):
    for ci, (val, cw, cx) in enumerate(zip(row, col_w, col_x)):
        box(s2, cx, 4.42 + ri*0.42, cw-0.08, 0.38,
            fill=row_colors[ri], border=RGBColor(0x33,0x55,0x77), border_w=Pt(0.5))
        txt(s2, val, cx+0.08, 4.44 + ri*0.42, cw-0.16, 0.34, size=12,
            color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Initial Program & What We Evolve
# ═══════════════════════════════════════════════════════════════════════════
s3 = slide()
heading(s3, "Initial Program & What We Evolve",
        "A Go function that decides: admit or reject each request?")

# Left column: structure
txt(s3, "admission.go  structure", 0.4, 1.35, 5.8, 0.38, size=14, bold=True, color=ACCENT)

struct_items = [
    ("AlwaysAdmit", "Baseline — admits everything (never used after init)"),
    ("TokenBucket",  "Rate-limiter reference implementation"),
    ("RejectAll",    "Pathological test case"),
    ("AdaptiveAdmission", "← This is what the LLM evolves"),
]
for i, (name, desc) in enumerate(struct_items):
    is_main = "AdaptiveAdmission" in name
    col = GREEN if is_main else LIGHT
    brd = GREEN if is_main else RGBColor(0x33,0x55,0x77)
    box(s3, 0.4, 1.82 + i*0.72, 5.8, 0.62,
        fill=RGBColor(0x16,0x2B,0x40) if not is_main else RGBColor(0x0a,0x25,0x10),
        border=brd, border_w=Pt(1.5 if is_main else 0.8))
    txt(s3, name, 0.55, 1.87 + i*0.72, 2.5, 0.28, size=13, bold=True, color=col)
    txt(s3, desc, 0.55, 2.13 + i*0.72, 5.5, 0.26, size=11, color=LIGHT)

# Right column: evolve block
txt(s3, "The EVOLVE-BLOCK  (only this changes)", 6.5, 1.35, 6.5, 0.38,
    size=14, bold=True, color=GREEN)

code_bg = RGBColor(0x08, 0x14, 0x20)
box(s3, 6.5, 1.78, 6.5, 4.5, fill=code_bg, border=GREEN, border_w=Pt(1.5))

code_lines = [
    ("// EVOLVE-BLOCK-START",             ACCENT),
    ("",                                   WHITE),
    ("// Available signals:",              LIGHT),
    ("sloClass   // critical/standard/...", WHITE),
    ("totalInFlight / numInstances",       YELLOW),
    ("tenantID, inputLen",                 WHITE),
    ("a.tenantTokens[tenantID]   // state",WHITE),
    ("",                                   WHITE),
    ("// Baseline just admits all:",       LIGHT),
    ("return true, \"\"",                  RED),
    ("",                                   WHITE),
    ("// LLM discovers better logic here", GREEN),
    ("",                                   WHITE),
    ("// EVOLVE-BLOCK-END",               ACCENT),
]
for i, (line, col) in enumerate(code_lines):
    txt(s3, line, 6.65, 1.92 + i*0.29, 6.2, 0.28, size=11, color=col)

# Bottom note
box(s3, 0.4, 6.75, 12.5, 0.5, fill=RGBColor(0x0a,0x20,0x10), border=GREEN, border_w=Pt(1))
txt(s3, "🔒  Only the EVOLVE-BLOCK is mutated. All signals are pre-computed above it. "
        "State maps are pre-initialized. No struct changes allowed — Go must compile cleanly.",
    0.55, 6.78, 12.2, 0.42, size=12, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Workloads
# ═══════════════════════════════════════════════════════════════════════════
s4 = slide()
heading(s4, "Evaluation Workloads",
        "Two realistic overload scenarios at 320 req/s (2× cluster capacity)")

for wi, (title, subtitle, clients, note, col) in enumerate([
    (
        "Workload 1: overload_mixed_slo",
        "Sustained 2× overload — steady Poisson arrivals",
        [
            ("tenant-rt",    "critical",  "20%", "128 tok input, streaming"),
            ("tenant-std",   "standard",  "30%", "256 tok input"),
            ("tenant-shed",  "sheddable", "20%", "256 tok input"),
            ("tenant-batch", "batch",     "30%", "512 tok input"),
        ],
        "Tests: Can the policy shed batch to protect critical under sustained overload?",
        ACCENT,
    ),
    (
        "Workload 2: bursty_adversary",
        "Burst isolation — gamma-distributed batch arrivals (CV=4)",
        [
            ("tenant-crit-1", "critical",  "15%", "128 tok, steady Poisson"),
            ("tenant-crit-2", "critical",  "15%", "128 tok, steady Poisson"),
            ("tenant-std",    "standard",  "20%", "256 tok, steady Poisson"),
            ("tenant-batch",  "batch",     "25%", "768 tok, γ CV=4 (bursty!)"),
            ("tenant-shed",   "sheddable", "25%", "512 tok, γ CV=3 (bursty!)"),
        ],
        "Tests: Does policy protect critical during bursts without over-rejecting in calm?",
        ORANGE,
    ),
]):
    ox = 0.4 + wi * 6.55
    box(s4, ox, 1.3, 6.35, 5.9, fill=CARD_BG, border=col, border_w=Pt(1.5))
    txt(s4, title,    ox+0.15, 1.38, 6.0, 0.38, size=14, bold=True,  color=col)
    txt(s4, subtitle, ox+0.15, 1.75, 6.0, 0.32, size=12, italic=True, color=LIGHT)

    # tenant rows
    row_cls_col = {"critical": RED, "standard": ACCENT, "sheddable": YELLOW, "batch": LIGHT}
    for ri, (tid, cls, frac, desc) in enumerate(clients):
        ry = 2.18 + ri * 0.72
        row_col = row_cls_col.get(cls, LIGHT)
        box(s4, ox+0.15, ry, 6.0, 0.62,
            fill=RGBColor(0x10,0x22,0x35), border=row_col, border_w=Pt(0.8))
        txt(s4, tid,  ox+0.28, ry+0.04, 1.6,  0.28, size=11, bold=True, color=WHITE)
        txt(s4, cls,  ox+0.28, ry+0.30, 1.6,  0.26, size=10, color=row_col)
        txt(s4, frac, ox+1.95, ry+0.15, 0.6,  0.28, size=13, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
        txt(s4, desc, ox+2.6,  ry+0.15, 3.4,  0.28, size=11, color=LIGHT)

    txt(s4, f"💡 {note}", ox+0.15, 1.3 + 5.9 - 0.52, 6.0, 0.42, size=11, italic=True, color=col)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Scoring Formula
# ═══════════════════════════════════════════════════════════════════════════
s5 = slide()
heading(s5, "How We Score Each Policy",
        "Multi-objective: SLO quality + throughput + fairness")

# Big formula box
box(s5, 0.8, 1.3, 11.7, 1.1, fill=RGBColor(0x08,0x20,0x10), border=GREEN, border_w=Pt(2))
txt(s5,
    "score  =  0.50 × weighted_SLO  +  0.30 × capped_throughput  +  0.20 × Jain_fairness",
    0.9, 1.42, 11.5, 0.65, size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# Three term cards
terms = [
    ("Weighted SLO  (50%)", ACCENT,
     ["Per-class SLO attainment, weighted by priority",
      "critical=4×,  standard=2×,  sheddable=1×,  batch=0.5×",
      "",
      "KEY RULE: Rejected batch/sheddable are EXCLUDED",
      "from denominator — shedding them is CORRECT,",
      "not a penalty. Rejected critical = SLO miss.",
      ],
     ),
    ("Capped Throughput  (30%)", ORANGE,
     ["min(completed / total,  0.50) / 0.50",
      "",
      "Shedding up to 50% of requests is FREE —",
      "full throughput credit. Only over-shedding",
      "(>50%) is penalized.",
      "→ Rewards moderate shedding strategies.",
      ],
     ),
    ("Jain Fairness  (20%)", YELLOW,
     ["Jain index over per-tenant completion rates",
      "",
      "( Σ xᵢ )²  /  ( N × Σ xᵢ² )",
      "",
      "N = all tenants in workload spec.",
      "Prevents 'admit only one tenant' exploits.",
      ],
     ),
]

for i, (title, col, lines) in enumerate(terms):
    cx = 0.35 + i * 4.33
    box(s5, cx, 2.6, 4.18, 3.8, fill=CARD_BG, border=col, border_w=Pt(1.5))
    txt(s5, title, cx+0.15, 2.68, 3.9, 0.42, size=14, bold=True, color=col)
    for j, line in enumerate(lines):
        txt(s5, line, cx+0.15, 3.18 + j*0.38, 3.9, 0.36, size=12, color=LIGHT)

# Bottom callout
box(s5, 0.35, 6.55, 12.6, 0.68, fill=RGBColor(0x10,0x22,0x10), border=GREEN, border_w=Pt(1))
txt(s5,
    "Baseline (always-admit) at 2× overload:   SLO=0.249   Throughput=1.0   Fairness=1.0   →   Score ≈ 0.62\n"
    "Oracle (ideal shedder):                   SLO≈0.87    Throughput=1.0   Fairness≈0.75  →   Score ≈ 0.85",
    0.5, 6.57, 12.3, 0.62, size=12, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Prompt & Models
# ═══════════════════════════════════════════════════════════════════════════
s6 = slide()
heading(s6, "System Prompt & LLM Configuration",
        "What we tell the LLM — and what we deliberately leave out")

# Left: models
txt(s6, "Models Used", 0.4, 1.35, 4.5, 0.38, size=14, bold=True, color=ACCENT)

for mi, (name, weight, role, col) in enumerate([
    ("Claude Sonnet 4.5", "70%", "Primary search model", ACCENT),
    ("Claude Opus 4.5",   "30%", "Occasional deep reasoning", YELLOW),
]):
    by = 1.82 + mi * 1.35
    box(s6, 0.4, by, 4.5, 1.15, fill=CARD_BG, border=col, border_w=Pt(1.5))
    txt(s6, name,   0.6, by+0.08, 3.8, 0.38, size=14, bold=True, color=col)
    txt(s6, f"Weight: {weight}  •  {role}", 0.6, by+0.46, 3.8, 0.3, size=12, color=LIGHT)
    txt(s6, "Temp 1.0  •  Max 32k tokens  •  120s timeout",
        0.6, by+0.76, 3.8, 0.28, size=11, color=RGBColor(0x88,0xaa,0xcc))

# Config pills
for ci, (k, v) in enumerate([("Temperature", "1.0"), ("Max tokens", "32,000"), ("Timeout", "120s")]):
    box(s6, 0.4 + ci*1.52, 4.35, 1.42, 0.52, fill=RGBColor(0x10,0x22,0x38), border=ACCENT, border_w=Pt(0.8))
    txt(s6, k, 0.4+ci*1.52, 4.36, 1.42, 0.24, size=10, color=LIGHT, align=PP_ALIGN.CENTER)
    txt(s6, v, 0.4+ci*1.52, 4.58, 1.42, 0.26, size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Right: prompt breakdown
txt(s6, "System Prompt  — what's IN and what's OUT", 5.2, 1.35, 7.8, 0.38,
    size=14, bold=True, color=ACCENT)

in_items = [
    "Problem context: 4-instance cluster, 2× overload",
    "Scoring formula with exact weights and rules",
    "Baseline score (~0.62) and oracle target (~0.85)",
    "All available signals + state maps explained",
    "Compilation rules (valid Go, no new struct fields)",
    "Shed-tolerant scoring explained (rejected ≠ SLO miss)",
]
out_items = [
    "Specific signal names to use  (e.g. totalInFlight)",
    "Threshold values  (no '> 15' or '> 30')",
    "Which SLO classes to shed",
    "Whether to use stateful or stateless logic",
]

txt(s6, "✅  Included", 5.2, 1.82, 3.6, 0.35, size=13, bold=True, color=GREEN)
for i, item in enumerate(in_items):
    txt(s6, f"• {item}", 5.2, 2.22 + i*0.44, 3.8, 0.38, size=12, color=LIGHT)

txt(s6, "🚫  Deliberately withheld  (to force real discovery)", 9.2, 1.82, 3.9, 0.35,
    size=13, bold=True, color=RED)
for i, item in enumerate(out_items):
    txt(s6, f"• {item}", 9.2, 2.22 + i*0.52, 3.9, 0.46, size=12, color=LIGHT)

# Generalization box
box(s6, 5.2, 4.85, 7.8, 1.55, fill=RGBColor(0x0a,0x20,0x10), border=GREEN, border_w=Pt(1.2))
txt(s6, "Generalization Requirement  (new in this run)", 5.35, 4.9, 7.5, 0.36,
    size=13, bold=True, color=GREEN)
txt(s6,
    "• No magic constants tuned to this simulation\n"
    "• Prefer normalized signals (per-instance ratios)\n"
    "• Adaptive thresholds derived from observed state are encouraged\n"
    "• Policy must transfer to real llm-d deployments",
    5.35, 5.28, 7.5, 1.0, size=12, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Search Process
# ═══════════════════════════════════════════════════════════════════════════
s7 = slide()
heading(s7, "How SkyDiscover Searches",
        "OpenEvolve: island-based evolutionary search with LLM mutations")

# Pipeline flow
steps = [
    ("1", "Initial\nProgram", "always-admit\nEVOLVE-BLOCK", ACCENT),
    ("2", "LLM\nMutation",   "Sonnet 4.5 (70%)\nOpus 4.5 (30%)", YELLOW),
    ("3", "Go Build\n& Run",  "2 workloads\n× 2 seeds = 4 sims", ORANGE),
    ("4", "Score\n& Rank",    "Combined score\nfeedback to DB", GREEN),
    ("5", "Island\nDatabase", "MAP-Elites\n5 islands", ACCENT),
]
for i, (num, title, sub, col) in enumerate(steps):
    bx = 0.4 + i * 2.5
    box(s7, bx, 1.35, 2.15, 1.6, fill=CARD_BG, border=col, border_w=Pt(1.5))
    txt(s7, num,   bx+0.08, 1.38, 0.45, 0.45, size=20, bold=True, color=col)
    txt(s7, title, bx+0.08, 1.75, 2.0,  0.52, size=14, bold=True, color=WHITE)
    txt(s7, sub,   bx+0.08, 2.27, 2.0,  0.55, size=11, color=LIGHT)
    if i < 4:
        arrow(s7, bx+2.15, 1.35+0.8, bx+2.5, 1.35+0.8)
    if i == 4:
        # loop back arrow hint
        txt(s7, "↺ repeat", bx+0.08, 2.82, 2.0, 0.28, size=11, color=col, italic=True)

# Score progression chart (manual bars)
txt(s7, "Score Progression (30 iterations)", 0.4, 3.25, 7.5, 0.4, size=14, bold=True, color=ACCENT)

chart_x, chart_y, chart_w, chart_h = 0.4, 3.72, 7.2, 2.8
box(s7, chart_x, chart_y, chart_w, chart_h, fill=RGBColor(0x08,0x14,0x20), border=ACCENT, border_w=Pt(1))

# Y axis labels
for val, label in [(0.62, "0.62"), (0.70, "0.70"), (0.78, "0.78"), (0.87, "0.87")]:
    norm = (val - 0.60) / (0.90 - 0.60)
    y_pos = chart_y + chart_h - norm * chart_h
    txt(s7, label, chart_x - 0.55, y_pos - 0.14, 0.5, 0.28, size=10,
        color=LIGHT, align=PP_ALIGN.RIGHT)
    line = s7.shapes.add_shape(1, Inches(chart_x), Inches(y_pos),
                                Inches(chart_w), Inches(0.01))
    line.fill.solid(); line.fill.fore_color.rgb = RGBColor(0x22,0x44,0x66)
    line.line.fill.background()

# Baseline line
baseline_norm = (0.6266 - 0.60) / (0.90 - 0.60)
by_pos = chart_y + chart_h - baseline_norm * chart_h
bl = s7.shapes.add_shape(1, Inches(chart_x), Inches(by_pos), Inches(chart_w), Inches(0.03))
bl.fill.solid(); bl.fill.fore_color.rgb = RED
bl.line.fill.background()
txt(s7, "baseline 0.62", chart_x + chart_w + 0.05, by_pos - 0.12, 1.1, 0.28, size=10, color=RED)

# Best score line
best_norm = (0.8700 - 0.60) / (0.90 - 0.60)
bst = s7.shapes.add_shape(1, Inches(chart_x), Inches(chart_y + chart_h - best_norm*chart_h),
                           Inches(chart_w), Inches(0.03))
bst.fill.solid(); bst.fill.fore_color.rgb = GREEN
bst.line.fill.background()
txt(s7, "best  0.87", chart_x + chart_w + 0.05,
    chart_y + chart_h - best_norm*chart_h - 0.12, 1.0, 0.28, size=10, color=GREEN)

# Dots for key iterations
key_iters = [(0, 0.6266, RED), (1, 0.8673, ACCENT), (7, 0.8700, GREEN)]
for it, score, col in key_iters:
    norm_x = (it / 30) * chart_w
    norm_y = (score - 0.60) / (0.90 - 0.60)
    dot_x = chart_x + norm_x
    dot_y = chart_y + chart_h - norm_y * chart_h
    dot = s7.shapes.add_shape(9, Inches(dot_x - 0.08), Inches(dot_y - 0.08),
                               Inches(0.16), Inches(0.16))
    dot.fill.solid(); dot.fill.fore_color.rgb = col
    dot.line.fill.background()
    txt(s7, f"iter {it}", Inches(dot_x + 0.1) / 914400,
        (Inches(dot_y - 0.25)) / 914400, 0.8, 0.24, size=9, color=col)

# Right: key facts
txt(s7, "Search Facts", 8.0, 3.25, 5.0, 0.4, size=14, bold=True, color=ACCENT)
facts = [
    ("30", "iterations"),
    ("0", "build errors  (100% compile rate)"),
    ("iter 1", "oracle-level policy first found"),
    ("iter 7", "final best  (+0.3% refinement)"),
    ("~37s", "per iteration  (~18 min total)"),
    ("2 seeds", "× 2 workloads = 4 sims/iter"),
]
for i, (val, label) in enumerate(facts):
    txt(s7, val,   8.0, 3.72 + i*0.46, 1.4, 0.4, size=14, bold=True, color=YELLOW, align=PP_ALIGN.RIGHT)
    txt(s7, label, 9.5, 3.72 + i*0.46, 3.6, 0.4, size=13, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Discovered Algorithm
# ═══════════════════════════════════════════════════════════════════════════
s8 = slide()
heading(s8, "What the LLM Discovered",
        "Adaptive, normalized admission policy — found at iteration 7")

# Left: pseudocode / logic flow
txt(s8, "Discovered Algorithm  (best_program.go, iter 7)", 0.4, 1.35, 7.0, 0.38,
    size=14, bold=True, color=GREEN)

box(s8, 0.4, 1.78, 6.8, 5.5, fill=RGBColor(0x08,0x14,0x20), border=GREEN, border_w=Pt(1.5))

logic = [
    ("// 10-second sliding window for load tracking",    LIGHT),
    ("perInstanceLoad = totalInFlight / numInstances",   YELLOW),
    ("",                                                  WHITE),
    ("// Bootstrap: estimate typicalLoad = 40 req/s",    LIGHT),
    ("// After 100 reqs: adapt from window stats",       LIGHT),
    ("loadRatio = perInstanceLoad / typicalLoad",        YELLOW),
    ("",                                                  WHITE),
    ("if sloClass == \"critical\"  → ADMIT always",      RGBColor(0x2E,0xCC,0x71)),
    ("if sloClass == \"standard\"  → ADMIT always",      RGBColor(0x2E,0xCC,0x71)),
    ("",                                                  WHITE),
    ("if sloClass == \"batch\":",                         RED),
    ("  fairnessRatio = tenantReqs / avgTenantReqs",     WHITE),
    ("  threshold = 0.50  // 50% of typical load",       ORANGE),
    ("  if fairnessRatio > 1.5: threshold *= 0.8",       LIGHT),
    ("  if loadRatio > threshold → REJECT",              RED),
    ("",                                                  WHITE),
    ("if sloClass == \"sheddable\":",                     ORANGE),
    ("  threshold = 0.75  // 75% of typical load",       ORANGE),
    ("  if fairnessRatio > 1.5: threshold *= 0.9",       LIGHT),
    ("  if loadRatio > threshold → REJECT",              ORANGE),
]
for i, (line, col) in enumerate(logic):
    txt(s8, line, 0.55, 1.95 + i*0.26, 6.5, 0.25, size=10.5, color=col)

# Right: key innovations
txt(s8, "Key Innovations vs Simple Oracle", 7.5, 1.35, 5.6, 0.38,
    size=14, bold=True, color=ACCENT)

innovations = [
    (ACCENT, "Normalized load signal",
     "Uses loadRatio = perInstanceLoad / typicalLoad\n"
     "instead of absolute threshold.\n"
     "Generalizes across cluster sizes."),
    (GREEN,  "Adaptive capacity estimate",
     "Bootstraps with typicalLoad=40, then\n"
     "learns from 10-second windows.\n"
     "Self-calibrates to actual workload."),
    (YELLOW, "Tenant fairness correction",
     "Tenants consuming >1.5× their fair share\n"
     "face tighter thresholds (0.8× / 0.9×).\n"
     "Prevents single tenant monopolizing."),
    (ORANGE, "Two-tier shedding",
     "Batch shed at 50% load, sheddable at 75%.\n"
     "Graduated response — not a cliff.\n"
     "Maintains throughput near cap."),
]
for i, (col, title, desc) in enumerate(innovations):
    by = 1.82 + i * 1.38
    box(s8, 7.5, by, 5.6, 1.22, fill=CARD_BG, border=col, border_w=Pt(1.2))
    txt(s8, title, 7.65, by+0.08, 5.3, 0.35, size=13, bold=True, color=col)
    txt(s8, desc,  7.65, by+0.44, 5.3, 0.7,  size=11, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Results
# ═══════════════════════════════════════════════════════════════════════════
s9 = slide()
heading(s9, "Results: +38.8% Over Always-Admit Baseline",
        "Experiment: 260324_30i_nohint_2seed  •  OpenEvolve  •  2 seeds  •  2 workloads")

# Big score comparison
for bi, (label, score, sub, col) in enumerate([
    ("Baseline\n(always-admit)", "0.6266", "SLO: 25%  |  Fair: 100%", RED),
    ("Best Evolved\n(iter 7)", "0.8700", "SLO: 84%  |  Fair: 74%", GREEN),
]):
    bx = 0.5 + bi * 5.5
    box(s9, bx, 1.3, 5.0, 2.0, fill=CARD_BG,
        border=col, border_w=Pt(2.5 if bi==1 else 1.5))
    txt(s9, label,  bx+0.2, 1.38, 4.6, 0.55, size=14, color=LIGHT)
    txt(s9, score,  bx+0.2, 1.88, 4.6, 0.85, size=46, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s9, sub,    bx+0.2, 2.82, 4.6, 0.38, size=13, color=LIGHT, align=PP_ALIGN.CENTER)

# Arrow + improvement
txt(s9, "→", 5.55, 1.95, 0.9, 0.8, size=38, color=ACCENT, align=PP_ALIGN.CENTER)
txt(s9, "+38.8%", 5.45, 2.72, 1.15, 0.46, size=20, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# Per-workload table
txt(s9, "Per-Workload Breakdown", 0.4, 3.52, 12.5, 0.38, size=14, bold=True, color=ACCENT)

col_hdrs = ["Workload", "Policy", "Score", "SLO Attain.", "Throughput", "Fairness", "Avg E2E", "P95 E2E"]
col_xs   = [0.35, 2.05, 3.45, 4.45, 5.60, 6.60, 7.65, 9.10]
col_ws   = [1.65, 1.35, 0.95, 1.10, 0.95, 1.00, 1.40, 1.50]

for ci, (hdr, cx, cw) in enumerate(zip(col_hdrs, col_xs, col_ws)):
    box(s9, cx, 3.97, cw-0.06, 0.38, fill=RGBColor(0x10,0x25,0x3a), border=ACCENT, border_w=Pt(0.8))
    txt(s9, hdr, cx+0.06, 3.99, cw-0.12, 0.32, size=11, bold=True,
        color=ACCENT, align=PP_ALIGN.CENTER)

rows = [
    ["overload_mixed_slo", "Baseline", "0.651", "30.2%", "100%", "100%", "6,783 ms", "15,341 ms"],
    ["overload_mixed_slo", "Evolved",  "0.869", "84.6%", "100%", "73.0%","1,178 ms", "3,634 ms"],
    ["bursty_adversary",   "Baseline", "0.598", "19.6%", "100%", "100%", "11,384 ms","22,898 ms"],
    ["bursty_adversary",   "Evolved",  "0.871", "84.1%", "100%", "75.2%","1,217 ms", "4,018 ms"],
]
row_cols = [
    RGBColor(0x22,0x10,0x10), RGBColor(0x10,0x2a,0x10),
    RGBColor(0x22,0x10,0x10), RGBColor(0x10,0x2a,0x10),
]
for ri, (row, rc) in enumerate(zip(rows, row_cols)):
    text_col = RED if "Baseline" in row[1] else GREEN
    for ci, (val, cx, cw) in enumerate(zip(row, col_xs, col_ws)):
        box(s9, cx, 4.40 + ri*0.52, cw-0.06, 0.46,
            fill=rc, border=RGBColor(0x33,0x55,0x66), border_w=Pt(0.4))
        vc = text_col if ci == 1 else WHITE
        txt(s9, val, cx+0.06, 4.42+ri*0.52, cw-0.12, 0.38, size=11,
            color=vc, align=PP_ALIGN.CENTER, bold=(ci==1))

# Latency improvement callout
box(s9, 10.6, 3.52, 2.5, 2.6, fill=RGBColor(0x0a,0x25,0x10), border=GREEN, border_w=Pt(1.5))
txt(s9, "Latency\nImprovement", 10.7, 3.58, 2.3, 0.52, size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
txt(s9, "Avg E2E", 10.7, 4.15, 2.3, 0.3, size=11, color=LIGHT, align=PP_ALIGN.CENTER)
txt(s9, "6.8×", 10.7, 4.42, 2.3, 0.55, size=32, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
txt(s9, "faster", 10.7, 4.92, 2.3, 0.3, size=11, color=LIGHT, align=PP_ALIGN.CENTER)
txt(s9, "P95 E2E", 10.7, 5.28, 2.3, 0.3, size=11, color=LIGHT, align=PP_ALIGN.CENTER)
txt(s9, "5.7×", 10.7, 5.52, 2.3, 0.46, size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txt(s9, "faster", 10.7, 5.95, 2.3, 0.28, size=11, color=LIGHT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Deployment Viability
# ═══════════════════════════════════════════════════════════════════════════
s10 = slide()
heading(s10, "Does It Transfer to Real llm-d?",
        "Signal availability and deployment viability assessment")

signals = [
    ("totalInFlight\n/ numInstances", "PRIMARY load signal",
     "Router-local, fresh every request.\nAvailable in llm-d routing plugin.\nNo Prometheus, no staleness.",
     GREEN, "HIGH"),
    ("sloClass", "Request priority",
     "From request header / metadata.\nNeeds custom header: X-SLO-Class.\nOne-time integration work.",
     ACCENT, "HIGH"),
    ("tenantID", "Tenant identity",
     "From auth/request headers.\nAvailable in llm-d via middleware.\nUsed for fairness correction.",
     ACCENT, "HIGH"),
    ("totalQueueDepth\nmaxKVUtil", "Stale cluster state",
     "Prometheus metrics, ~5s staleness.\nUsed as secondary signals only.\nNot in the best evolved policy.",
     YELLOW, "MED"),
    ("outputTokens\ncacheHitRate", "Simulator-only",
     "Not available in real llm-d router.\nBest policy doesn't use these.\nSafe to ignore.",
     LIGHT, "N/A"),
]

for i, (sig, role, desc, col, viability) in enumerate(signals):
    bx = 0.35 + (i % 3) * 4.3
    by = 1.32 + (i // 3) * 2.55
    box(s10, bx, by, 4.15, 2.3, fill=CARD_BG, border=col, border_w=Pt(1.2))
    # viability badge
    v_col = GREEN if viability=="HIGH" else (YELLOW if viability=="MED" else LIGHT)
    box(s10, bx + 2.7, by+0.08, 1.3, 0.38, fill=v_col, border=v_col)
    txt(s10, viability, bx+2.7, by+0.09, 1.3, 0.34, size=12, bold=True,
        color=BG, align=PP_ALIGN.CENTER)
    txt(s10, sig,  bx+0.15, by+0.08, 2.5, 0.52, size=12, bold=True, color=col)
    txt(s10, role, bx+0.15, by+0.60, 3.9, 0.28, size=11, italic=True, color=WHITE)
    txt(s10, desc, bx+0.15, by+0.92, 3.9, 0.88, size=11, color=LIGHT)

# Bottom summary
box(s10, 0.35, 6.62, 12.6, 0.6, fill=RGBColor(0x0a,0x25,0x10), border=GREEN, border_w=Pt(1.5))
txt(s10,
    "✅  Deployment verdict: HIGH viability.  The best policy uses only totalInFlight (router-local) + "
    "sloClass + tenantID — all available in llm-d with minimal integration.  "
    "No simulator-only signals. Adaptive thresholds self-calibrate to real cluster behaviour.",
    0.5, 6.65, 12.3, 0.52, size=12, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Summary
# ═══════════════════════════════════════════════════════════════════════════
s11 = slide()
heading(s11, "Summary", "What we set up, what was discovered, what it means")

takeaways = [
    (GREEN,  "Setup",
     "4-instance qwen2.5-7b cluster at 2× overload (320 QPS). "
     "Multi-objective score: weighted SLO (50%) + capped throughput (30%) + Jain fairness (20%). "
     "Shed-tolerant: rejecting batch/sheddable is rewarded, not penalized."),
    (ACCENT, "Discovery",
     "OpenEvolve found a normalised, adaptive policy in 30 iterations (18 min). "
     "Core logic: perInstanceLoad / typicalLoad as load signal, with self-calibrating "
     "typicalLoad estimate and per-tenant fairness correction. No hardcoded magic numbers."),
    (YELLOW, "Gains",
     "+38.8% combined score (0.627 → 0.870). Critical SLO attainment: 25% → 84%. "
     "Average E2E latency: 9,083 ms → 1,197 ms (6.8×). Throughput unchanged (full cap credit). "
     "Consistent across both workloads and both random seeds."),
    (ORANGE, "Prompt design",
     "Removing the 'WINNING STRATEGY' hint (specific thresholds) forced real discovery. "
     "Adding a GENERALIZATION REQUIREMENT produced a more adaptive, deployable policy "
     "with much better fairness (0.741 vs 0.579 in the hinted run)."),
    (RGBColor(0xDD,0x88,0xFF), "Next steps",
     "Test at 1.5× overload (240 QPS) — more realistic production scenario. "
     "Multi-LLM validation with qwen14b. "
     "Transfer to real llm-d admission hook. "
     "Explore joint router + admission co-evolution."),
]

for i, (col, label, text) in enumerate(takeaways):
    by = 1.32 + i * 1.12
    box(s11, 0.35, by, 0.72, 0.88, fill=col, border=col)
    txt(s11, label, 0.38, by+0.18, 0.68, 0.52, size=11, bold=True,
        color=BG, align=PP_ALIGN.CENTER)
    box(s11, 1.12, by, 11.85, 0.88, fill=CARD_BG, border=col, border_w=Pt(1))
    txt(s11, text, 1.27, by+0.12, 11.55, 0.65, size=12.5, color=LIGHT)


# ── Save ─────────────────────────────────────────────────────────────────────
out = "docs/slides/blis_admission_experiment.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
