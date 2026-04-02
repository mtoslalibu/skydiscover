"""
BLIS Router — SkyDiscover Experiment Deck
Run: python3 docs/slides/blis_router_deck.py
Output: docs/slides/blis_router_experiment.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0D, 0x1B, 0x2A)   # dark navy
ACCENT    = RGBColor(0x00, 0xC2, 0xFF)   # cyan
GREEN     = RGBColor(0x2E, 0xCC, 0x71)   # success green
ORANGE    = RGBColor(0xFF, 0x8C, 0x00)   # warning orange
RED       = RGBColor(0xFF, 0x4D, 0x4D)   # danger red
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCC, 0xDD, 0xEE)
CARD_BG   = RGBColor(0x16, 0x2B, 0x40)
YELLOW    = RGBColor(0xFF, 0xD7, 0x00)

W, H = Inches(13.33), Inches(7.5)   # 16:9 widescreen

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ── Helpers ──────────────────────────────────────────────────────────────────

def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return s

def box(s, x, y, w, h, fill=None, border=None, border_w=Pt(1.5), radius=None):
    shape = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background() if border is None else None
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = border_w
    return shape

def txt(s, text, x, y, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def heading(s, title, subtitle=None):
    txt(s, title, 0.4, 0.18, 12.5, 0.7, size=32, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(s, subtitle, 0.4, 0.82, 12.5, 0.4, size=16, color=LIGHT, align=PP_ALIGN.LEFT)
    line = s.shapes.add_shape(1, Inches(0.4), Inches(1.15), Inches(12.5), Inches(0.03))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

def card(s, x, y, w, h, title, body_lines, title_color=ACCENT, body_size=14):
    box(s, x, y, w, h, fill=CARD_BG, border=ACCENT, border_w=Pt(1))
    txt(s, title, x+0.15, y+0.12, w-0.3, 0.38, size=15, bold=True, color=title_color)
    body = "\n".join(body_lines)
    txt(s, body, x+0.15, y+0.52, w-0.3, h-0.65, size=body_size, color=LIGHT)

def bullet_block(s, x, y, w, lines, size=15, gap=0.34):
    for i, (bullet, text) in enumerate(lines):
        txt(s, bullet, x, y + i*gap, 0.3, gap, size=size, color=ACCENT, bold=True)
        txt(s, text,   x+0.28, y + i*gap, w-0.28, gap, size=size, color=WHITE)

def arrow(s, x1, y1, x2, y2, color=ACCENT):
    connector = s.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(2)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════
s1 = slide()

txt(s1, "BLIS Router", 1.0, 1.6, 11.3, 1.1,
    size=48, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txt(s1, "Discovering adaptive routing algorithms for LLM inference clusters",
    1.0, 2.75, 11.3, 0.6, size=22, color=LIGHT, align=PP_ALIGN.CENTER)
txt(s1, "using SkyDiscover + OpenEvolve",
    1.0, 3.3, 11.3, 0.5, size=18, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

line = s1.shapes.add_shape(1, Inches(3.5), Inches(4.0), Inches(6.3), Inches(0.04))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

txt(s1, "Experiment:  260312_50i_openevolve_v2wl  •  March 2026",
    1.0, 4.2, 11.3, 0.4, size=14, color=LIGHT, align=PP_ALIGN.CENTER)

for i, (val, label, col) in enumerate([
    ("+11.5%", "vs best hand-tuned baseline", GREEN),
    ("+42.0%", "vs LLQ (common default)", ACCENT),
    ("50 iters", "OpenEvolve search", ORANGE),
]):
    bx = 1.8 + i * 3.3
    box(s1, bx, 5.1, 2.9, 1.3, fill=CARD_BG, border=col)
    txt(s1, val,   bx, 5.2,  2.9, 0.55, size=28, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s1, label, bx, 5.75, 2.9, 0.4,  size=13, color=LIGHT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Setup
# ═══════════════════════════════════════════════════════════════════════════
s2 = slide()
heading(s2, "The Problem: Routing in an LLM Inference Cluster",
        "How do you route requests across instances to minimize latency?")

# Left: cluster diagram
txt(s2, "4-Instance Cluster  (qwen2.5-7b / H100)", 0.4, 1.35, 5.5, 0.4,
    size=14, bold=True, color=ACCENT)

for i in range(4):
    bx, by = 0.5 + i*1.22, 1.85
    box(s2, bx, by, 1.05, 0.8, fill=RGBColor(0x1a,0x3a,0x5c), border=ACCENT, border_w=Pt(1))
    txt(s2, f"GPU {i+1}", bx, by+0.05, 1.05, 0.35, size=13, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s2, "H100", bx, by+0.4, 1.05, 0.28, size=11, color=LIGHT, align=PP_ALIGN.CENTER)

# Router box above
box(s2, 1.1, 1.35, 2.9, 0.42, fill=RGBColor(0x0a,0x25,0x3a), border=ACCENT, border_w=Pt(1.5))
txt(s2, "llm-d  Router", 1.1, 1.38, 2.9, 0.36, size=13, bold=True,
    color=ACCENT, align=PP_ALIGN.CENTER)

# KV cache note
box(s2, 0.4, 2.85, 5.2, 0.45, fill=RGBColor(0x0a,0x1a,0x2a), border=ACCENT)
txt(s2, "Each GPU has KV cache  •  Prefix hits = big speedup",
    0.55, 2.88, 5.0, 0.38, size=13, color=ACCENT)

# Right: what goes wrong
txt(s2, "Why Routing Matters", 6.2, 1.35, 6.5, 0.4,
    size=14, bold=True, color=YELLOW)

problems = [
    ("🔴", "Wrong instance: no prefix cache hit — pay full prefill cost"),
    ("🔴", "Hot instance: high InFlightRequests causes queuing"),
    ("🔴", "KV overflow: memory pressure triggers preemption"),
    ("🟡", "Naïve load balance ignores prefix locality entirely"),
]
for i, (icon, text) in enumerate(problems):
    txt(s2, icon, 6.2, 1.85 + i*0.52, 0.4, 0.45, size=18)
    txt(s2, text, 6.65, 1.85 + i*0.52, 6.3, 0.45, size=14, color=LIGHT)

# Bottom: signal freshness overview
txt(s2, "Available Signals  (freshness matters!)", 0.4, 3.55, 12.5, 0.38,
    size=14, bold=True, color=ACCENT)

sig_rows = [
    ("InFlightRequests", "FRESH — every request", "Router-local counter, no Prometheus lag", GREEN),
    ("KVUtilization / FreeKVBlocks", "STALE — ~5s", "Prometheus scrape, may lag under burst", YELLOW),
    ("QueueDepth / BatchSize", "STALE — ~5s", "Prometheus scrape", YELLOW),
    ("Prefix-affinity score", "FRESH — every request", "Router-local LRU cache index", GREEN),
]
col_ws = [2.8, 2.2, 4.8]
col_xs = [0.4, 3.25, 5.5]
col_hdrs = ["Signal", "Freshness", "Source"]
for ci, (hdr, cx, cw) in enumerate(zip(col_hdrs, col_xs, col_ws)):
    box(s2, cx, 4.0, cw-0.08, 0.36, fill=RGBColor(0x10,0x25,0x3a), border=ACCENT, border_w=Pt(0.8))
    txt(s2, hdr, cx+0.08, 4.02, cw-0.16, 0.30, size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
for ri, (sig, fresh, src, col) in enumerate(sig_rows):
    vals = [sig, fresh, src]
    for ci, (val, cx, cw) in enumerate(zip(vals, col_xs, col_ws)):
        box(s2, cx, 4.4 + ri*0.42, cw-0.08, 0.38,
            fill=CARD_BG, border=RGBColor(0x33,0x55,0x77), border_w=Pt(0.5))
        vc = col if ci == 1 else (ACCENT if ci == 0 else LIGHT)
        txt(s2, val, cx+0.08, 4.42 + ri*0.42, cw-0.16, 0.34, size=11,
            color=vc, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Initial Program & What We Evolve
# ═══════════════════════════════════════════════════════════════════════════
s3 = slide()
heading(s3, "Initial Program & What We Evolve",
        "A Go function inside the router that scores every instance for each request")

# Left column: routing.go structure
txt(s3, "routing.go  structure", 0.4, 1.35, 5.8, 0.38, size=14, bold=True, color=ACCENT)

struct_items = [
    ("RoundRobin",      "Uniform round-robin, no load awareness"),
    ("LeastLoaded",     "Minimum InFlightRequests — ignores prefix"),
    ("AlwaysBusiest",   "Pathological test case"),
    ("WeightedScoring", "← This is what the LLM evolves"),
]
for i, (name, desc) in enumerate(struct_items):
    is_main = "WeightedScoring" in name
    col = GREEN if is_main else LIGHT
    brd = GREEN if is_main else RGBColor(0x33,0x55,0x77)
    box(s3, 0.4, 1.82 + i*0.72, 5.8, 0.62,
        fill=RGBColor(0x16,0x2B,0x40) if not is_main else RGBColor(0x0a,0x25,0x10),
        border=brd, border_w=Pt(1.5 if is_main else 0.8))
    txt(s3, name, 0.55, 1.87 + i*0.72, 2.5, 0.28, size=13, bold=True, color=col)
    txt(s3, desc, 0.55, 2.13 + i*0.72, 5.5, 0.26, size=11, color=LIGHT)

# Right column: evolve block
txt(s3, "The EVOLVE-BLOCK  (only this changes)", 6.5, 1.35, 6.5, 0.38,
    size=14, bold=True, color=GREEN)

code_bg = RGBColor(0x08, 0x14, 0x20)
box(s3, 6.5, 1.78, 6.5, 4.5, fill=code_bg, border=GREEN, border_w=Pt(1.5))

code_lines = [
    ("// EVOLVE-BLOCK-START",              ACCENT),
    ("",                                    WHITE),
    ("// Available per-instance:",          LIGHT),
    ("snap.InFlightRequests   // FRESH",   YELLOW),
    ("snap.KVUtilization      // stale",   WHITE),
    ("snap.FreeKVBlocks        // stale",  WHITE),
    ("allDimScores[i][snap.ID]  // scorer", WHITE),
    ("",                                    WHITE),
    ("// Baseline: pure weighted sum",      LIGHT),
    ("score += dimScore * weight",         RED),
    ("",                                    WHITE),
    ("// LLM discovers adaptive logic",    GREEN),
    ("// e.g. decay prefix when hot...",   GREEN),
    ("",                                    WHITE),
    ("// EVOLVE-BLOCK-END",               ACCENT),
]
for i, (line, col) in enumerate(code_lines):
    txt(s3, line, 6.65, 1.92 + i*0.29, 6.2, 0.28, size=11, color=col)

# Bottom note
box(s3, 0.4, 6.75, 12.5, 0.5, fill=RGBColor(0x0a,0x20,0x10), border=GREEN, border_w=Pt(1))
txt(s3, "🔒  Only the EVOLVE-BLOCK is mutated. Scorer pipeline, state, and interface stay fixed. "
        "Go must compile cleanly. The router is drop-in-replaceable in real llm-d.",
    0.55, 6.78, 12.2, 0.42, size=12, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Workloads
# ═══════════════════════════════════════════════════════════════════════════
s4 = slide()
heading(s4, "Evaluation Workloads",
        "Two realistic traffic profiles — general load and prefix-heavy")

for wi, (title, subtitle, clients, note, col) in enumerate([
    (
        "Workload 1: glia_40qps",
        "General bursty traffic — 40 QPS, 1000 requests",
        [
            ("ShareGPT-like (90%)", "~500 tok input, ~250 tok output", "Gamma CV=7.3 (bursty)"),
            ("Heavy prompt (5%)",   "~5000 tok input, ~250 tok output", "10× longer prompts"),
            ("Heavy decode (5%)",   "~500 tok input, ~2500 tok output", "10× longer decodes"),
        ],
        "Tests: Load balancing under bursty arrivals. No prefix sharing.",
        ACCENT,
    ),
    (
        "Workload 2: glia_prefix_heavy",
        "Prefix-heavy traffic — 85 QPS, 1500 requests, 6 prefix groups",
        [
            ("Group A (45%)", "14,336-token shared prefix", "Dominant group"),
            ("Groups B–E (8–18% each)", "14,336-token shared prefixes", "Medium groups"),
            ("Group F (7%)", "No shared prefix", "No cache benefit"),
        ],
        "Tests: Prefix locality vs load balance. Hotspotting risk if routing ignores load.",
        ORANGE,
    ),
]):
    ox = 0.4 + wi * 6.55
    box(s4, ox, 1.3, 6.35, 5.9, fill=CARD_BG, border=col, border_w=Pt(1.5))
    txt(s4, title,    ox+0.15, 1.38, 6.0, 0.38, size=14, bold=True,  color=col)
    txt(s4, subtitle, ox+0.15, 1.75, 6.0, 0.32, size=12, italic=True, color=LIGHT)

    row_cols_list = [ACCENT, YELLOW, ORANGE]
    for ri, (ctype, desc, detail) in enumerate(clients):
        ry = 2.18 + ri * 0.9
        row_col = row_cols_list[ri % len(row_cols_list)]
        box(s4, ox+0.15, ry, 6.0, 0.78,
            fill=RGBColor(0x10,0x22,0x35), border=row_col, border_w=Pt(0.8))
        txt(s4, ctype,  ox+0.28, ry+0.05, 2.4, 0.28, size=12, bold=True, color=WHITE)
        txt(s4, desc,   ox+0.28, ry+0.33, 3.4, 0.24, size=11, color=row_col)
        txt(s4, detail, ox+2.8,  ry+0.05, 3.2, 0.28, size=11, color=LIGHT, italic=True)

    txt(s4, f"💡 {note}", ox+0.15, 1.3 + 5.9 - 0.52, 6.0, 0.42, size=11, italic=True, color=col)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Scoring Formula
# ═══════════════════════════════════════════════════════════════════════════
s5 = slide()
heading(s5, "How We Score Each Router",
        "Minimize latency across both workloads and both random seeds")

# Big formula box
box(s5, 0.8, 1.3, 11.7, 1.1, fill=RGBColor(0x08,0x20,0x10), border=GREEN, border_w=Pt(2))
txt(s5,
    "score  =  0.5 × E2E_mean  +  0.5 × E2E_P95   (per workload, then averaged)",
    0.9, 1.42, 11.5, 0.65, size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# Three term cards
terms = [
    ("E2E Mean Latency  (50%)", ACCENT,
     ["Average end-to-end latency per request",
      "across all completed requests.",
      "",
      "Lower = better. Normalized per workload",
      "then combined as mean improvement.",
      "→ Rewards overall throughput quality.",
      ]),
    ("E2E P95 Latency  (50%)", YELLOW,
     ["95th percentile end-to-end latency.",
      "",
      "Lower = better. Captures tail behaviour",
      "— important for real-time and interactive",
      "SLO tiers that care about worst-case.",
      "→ Penalizes hotspot-induced outliers.",
      ]),
    ("Baseline Normalisation", ORANGE,
     ["All scores expressed as % improvement",
      "vs 1:1 (50/50 prefix-affinity + load-balance).",
      "",
      "1:1 is the 'sensible default' hand-tuned",
      "by engineers — a strong baseline.",
      "→ +11.5% means 11.5% less latency.",
      ]),
]

for i, (title, col, lines) in enumerate(terms):
    cx = 0.35 + i * 4.33
    box(s5, cx, 2.6, 4.18, 3.8, fill=CARD_BG, border=col, border_w=Pt(1.5))
    txt(s5, title, cx+0.15, 2.68, 3.9, 0.42, size=14, bold=True, color=col)
    for j, line in enumerate(lines):
        txt(s5, line, cx+0.15, 3.18 + j*0.38, 3.9, 0.36, size=12, color=LIGHT)

# Bottom callout: baselines ranked
box(s5, 0.35, 6.55, 12.6, 0.68, fill=RGBColor(0x10,0x22,0x10), border=GREEN, border_w=Pt(1))
txt(s5,
    "Baselines ranked (vs 1:1):   LLQ = −61.0%   |   Glia = −14.3%   |   3:2:2 = −0.7%   |   "
    "1:1 = 0% (control)   |   Evolved = +11.5%",
    0.5, 6.57, 12.3, 0.62, size=12, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Prompt & Models
# ═══════════════════════════════════════════════════════════════════════════
s6 = slide()
heading(s6, "System Prompt & LLM Configuration",
        "What we tell the LLM — and what we leave to discovery")

# Left: models
txt(s6, "Models Used", 0.4, 1.35, 4.5, 0.38, size=14, bold=True, color=ACCENT)

for mi, (name, weight, role, col) in enumerate([
    ("Claude Sonnet 4.5", "70%", "Primary search model", ACCENT),
    ("Claude Opus 4.5",   "30%", "Occasional deep reasoning", YELLOW),
]):
    by = 1.82 + mi * 1.35
    box(s6, 0.4, by, 4.5, 1.15, fill=CARD_BG, border=col, border_w=Pt(1.5))
    txt(s6, name,   0.6, by+0.08, 3.8, 0.38, size=14, bold=True, color=col)
    txt(s6, f"Weight: {weight}  •  {role}", 0.6, by+0.46, 3.8, 0.3, size=12, color=LIGHT)
    txt(s6, "Temp 1.0  •  Max 32k tokens  •  120s timeout",
        0.6, by+0.76, 3.8, 0.28, size=11, color=RGBColor(0x88,0xaa,0xcc))

for ci, (k, v) in enumerate([("Temperature", "1.0"), ("Max tokens", "32,000"), ("Timeout", "120s")]):
    box(s6, 0.4 + ci*1.52, 4.35, 1.42, 0.52, fill=RGBColor(0x10,0x22,0x38), border=ACCENT, border_w=Pt(0.8))
    txt(s6, k, 0.4+ci*1.52, 4.36, 1.42, 0.24, size=10, color=LIGHT, align=PP_ALIGN.CENTER)
    txt(s6, v, 0.4+ci*1.52, 4.58, 1.42, 0.26, size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Right: prompt breakdown
txt(s6, "System Prompt  — what's IN and what's OUT", 5.2, 1.35, 7.8, 0.38,
    size=14, bold=True, color=ACCENT)

in_items = [
    "Cluster setup: 4 instances, model, hardware",
    "Available signals + their freshness (fresh vs ~5s stale)",
    "Scorer pipeline API (prefix-affinity, load-balance weights)",
    "Workload descriptions + scoring formula",
    "Baseline comparison table (LLQ, Glia, 1:1 scores)",
    "Compilation rules (valid Go, must compile cleanly)",
]
out_items = [
    "Which signals to combine or prioritize",
    "Decay functions or adaptive weights",
    "How to handle KV pressure",
    "Whether to use stateful logic",
]

txt(s6, "✅  Included", 5.2, 1.82, 3.6, 0.35, size=13, bold=True, color=GREEN)
for i, item in enumerate(in_items):
    txt(s6, f"• {item}", 5.2, 2.22 + i*0.44, 3.8, 0.38, size=12, color=LIGHT)

txt(s6, "🚫  Deliberately left to LLM", 9.2, 1.82, 3.9, 0.35,
    size=13, bold=True, color=RED)
for i, item in enumerate(out_items):
    txt(s6, f"• {item}", 9.2, 2.22 + i*0.52, 3.9, 0.46, size=12, color=LIGHT)

# Key insight box
box(s6, 5.2, 4.85, 7.8, 1.55, fill=RGBColor(0x0a,0x20,0x10), border=GREEN, border_w=Pt(1.2))
txt(s6, "Design Insight", 5.35, 4.9, 7.5, 0.36,
    size=13, bold=True, color=GREEN)
txt(s6,
    "• Freshness matters: InFlightRequests is synchronous, KVUtil is 5s stale\n"
    "• LLM must discover when to trust each signal\n"
    "• Tension: prefix locality vs hot-instance overload\n"
    "• No hints given — LLM must find the right balance",
    5.35, 5.28, 7.5, 1.0, size=12, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Search Process
# ═══════════════════════════════════════════════════════════════════════════
s7 = slide()
heading(s7, "How SkyDiscover Searches",
        "OpenEvolve: island-based evolutionary search with LLM mutations")

# Pipeline flow
steps = [
    ("1", "Initial\nProgram", "1:1 weighted\nEVOLVE-BLOCK", ACCENT),
    ("2", "LLM\nMutation",   "Sonnet 4.5 (70%)\nOpus 4.5 (30%)", YELLOW),
    ("3", "Go Build\n& Run",  "2 workloads\n× 2 seeds = 4 sims", ORANGE),
    ("4", "Score\n& Rank",    "E2E + P95\nvs 1:1 baseline", GREEN),
    ("5", "Island\nDatabase", "MAP-Elites\n5 islands", ACCENT),
]
for i, (num, title, sub, col) in enumerate(steps):
    bx = 0.4 + i * 2.5
    box(s7, bx, 1.35, 2.15, 1.6, fill=CARD_BG, border=col, border_w=Pt(1.5))
    txt(s7, num,   bx+0.08, 1.38, 0.45, 0.45, size=20, bold=True, color=col)
    txt(s7, title, bx+0.08, 1.75, 2.0,  0.52, size=14, bold=True, color=WHITE)
    txt(s7, sub,   bx+0.08, 2.27, 2.0,  0.55, size=11, color=LIGHT)
    if i < 4:
        arrow(s7, bx+2.15, 1.35+0.8, bx+2.5, 1.35+0.8)
    if i == 4:
        txt(s7, "↺ repeat", bx+0.08, 2.82, 2.0, 0.28, size=11, color=col, italic=True)

# Score progression chart
txt(s7, "Score Progression (50 iterations, % improvement vs 1:1)", 0.4, 3.25, 7.5, 0.4,
    size=14, bold=True, color=ACCENT)

chart_x, chart_y, chart_w, chart_h = 0.4, 3.72, 7.2, 2.8
box(s7, chart_x, chart_y, chart_w, chart_h, fill=RGBColor(0x08,0x14,0x20), border=ACCENT, border_w=Pt(1))

# Y-axis grid lines at 0%, 4%, 8%, 11.5%
for val, label in [(0, "0%"), (4, "4%"), (8, "8%"), (11.5, "11.5%")]:
    norm = val / 13.0
    y_pos = chart_y + chart_h - norm * chart_h
    txt(s7, label, chart_x - 0.65, y_pos - 0.14, 0.6, 0.28, size=10,
        color=LIGHT, align=PP_ALIGN.RIGHT)
    gline = s7.shapes.add_shape(1, Inches(chart_x), Inches(y_pos),
                                Inches(chart_w), Inches(0.01))
    gline.fill.solid(); gline.fill.fore_color.rgb = RGBColor(0x22,0x44,0x66)
    gline.line.fill.background()

# Baseline (1:1) line at 0
baseline_norm = 0.0
by_pos = chart_y + chart_h - baseline_norm * chart_h
bl = s7.shapes.add_shape(1, Inches(chart_x), Inches(by_pos), Inches(chart_w), Inches(0.03))
bl.fill.solid(); bl.fill.fore_color.rgb = RED
bl.line.fill.background()
txt(s7, "1:1 baseline  0%", chart_x + chart_w + 0.05, by_pos - 0.12, 1.4, 0.28, size=10, color=RED)

# Best score line
best_norm = 11.46 / 13.0
bst_y = chart_y + chart_h - best_norm * chart_h
bst = s7.shapes.add_shape(1, Inches(chart_x), Inches(bst_y), Inches(chart_w), Inches(0.03))
bst.fill.solid(); bst.fill.fore_color.rgb = GREEN
bst.line.fill.background()
txt(s7, "best  +11.5%", chart_x + chart_w + 0.05, bst_y - 0.12, 1.2, 0.28, size=10, color=GREEN)

# Key iteration dots
key_iters = [(0, 0.0, RED), (1, 7.18, ACCENT), (2, 10.66, YELLOW), (16, 11.46, GREEN)]
for it, score, col in key_iters:
    norm_x = (it / 50) * chart_w
    norm_y = score / 13.0
    dot_x = chart_x + norm_x
    dot_y = chart_y + chart_h - norm_y * chart_h
    dot = s7.shapes.add_shape(9, Inches(dot_x - 0.08), Inches(dot_y - 0.08),
                               Inches(0.16), Inches(0.16))
    dot.fill.solid(); dot.fill.fore_color.rgb = col
    dot.line.fill.background()
    txt(s7, f"iter {it}", dot_x + 0.1, dot_y - 0.28, 0.9, 0.24, size=9, color=col)

# Right: key facts
txt(s7, "Search Facts", 8.0, 3.25, 5.0, 0.4, size=14, bold=True, color=ACCENT)
facts = [
    ("50", "iterations  (31.8 min total)"),
    ("0", "build errors  (100% compile rate)"),
    ("iter 1", "+7.2%  first attempt"),
    ("iter 2", "+10.7%  major gain"),
    ("iter 16", "final best  +11.5%"),
    ("~38s", "per iteration  (2 wl × 2 seeds)"),
]
for i, (val, label) in enumerate(facts):
    txt(s7, val,   8.0, 3.72 + i*0.46, 1.4, 0.4, size=14, bold=True, color=YELLOW, align=PP_ALIGN.RIGHT)
    txt(s7, label, 9.5, 3.72 + i*0.46, 3.6, 0.4, size=13, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Discovered Algorithm
# ═══════════════════════════════════════════════════════════════════════════
s8 = slide()
heading(s8, "What the LLM Discovered",
        "Adaptive prefix-affinity router — found at iteration 16")

# Left: pseudocode
txt(s8, "Discovered Algorithm  (best_program.go, iter 16)", 0.4, 1.35, 7.0, 0.38,
    size=14, bold=True, color=GREEN)

box(s8, 0.4, 1.78, 6.8, 5.5, fill=RGBColor(0x08,0x14,0x20), border=GREEN, border_w=Pt(1.5))

logic = [
    ("// Run all scorers (prefix-affinity, load-balance)", LIGHT),
    ("allDimScores[i] = scorer(req, snapshots)",           YELLOW),
    ("",                                                    WHITE),
    ("// Find best prefix-affinity instance + load",       LIGHT),
    ("bestPrefixID = argmax(allDimScores[0])",             YELLOW),
    ("cachedLoad = snapshots[bestPrefixID].InFlightRequests",WHITE),
    ("delta = cachedLoad - minInflight",                   WHITE),
    ("",                                                    WHITE),
    ("// Decay prefix weight when cached instance is hot", LIGHT),
    ("if delta > 0:",                                      ORANGE),
    ("  decay = 1 / (1 + 0.6 * delta)",                   ORANGE),
    ("  aw[prefix] = weight[prefix] * decay",              ORANGE),
    ("  aw[load]   = 1.0 - aw[prefix]",                   ORANGE),
    ("",                                                    WHITE),
    ("// KV pressure penalty",                             LIGHT),
    ("if KVUtil > 0.9:",                                   RED),
    ("  score -= 0.5 * (KVUtil - 0.9) / 0.1",             RED),
    ("",                                                    WHITE),
    ("// Fresh load tiebreaker",                           LIGHT),
    ("score += 0.01 / (1 + InFlightRequests)",             GREEN),
]
for i, (line, col) in enumerate(logic):
    txt(s8, line, 0.55, 1.95 + i*0.26, 6.5, 0.25, size=10.5, color=col)

# Right: key innovations
txt(s8, "Key Innovations vs 1:1 Baseline", 7.5, 1.35, 5.6, 0.38,
    size=14, bold=True, color=ACCENT)

innovations = [
    (ORANGE, "Adaptive prefix-affinity decay",
     "When the best-cache instance is overloaded,\n"
     "decay = 1/(1 + 0.6 × load_delta) reduces\n"
     "its weight. Prevents cache hotspotting."),
    (RED,    "KV pressure penalty",
     "Instance with KVUtil > 90% is penalized\n"
     "by −0.5 × (KVUtil − 0.9) / 0.1.\n"
     "Avoids memory pressure before preemption."),
    (GREEN,  "Fresh load tiebreaker",
     "Adds +0.01 / (1 + InFlightRequests).\n"
     "Uses the freshest signal (per-request)\n"
     "to break ties without waiting for Prometheus."),
    (YELLOW, "All signals are production-available",
     "InFlightRequests: router-local (fresh).\n"
     "KVUtilization: Prometheus (~5s stale).\n"
     "Prefix scores: router LRU cache (fresh)."),
]
for i, (col, title, desc) in enumerate(innovations):
    by = 1.82 + i * 1.38
    box(s8, 7.5, by, 5.6, 1.22, fill=CARD_BG, border=col, border_w=Pt(1.2))
    txt(s8, title, 7.65, by+0.08, 5.3, 0.35, size=13, bold=True, color=col)
    txt(s8, desc,  7.65, by+0.44, 5.3, 0.7,  size=11, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Results vs All Baselines
# ═══════════════════════════════════════════════════════════════════════════
s9 = slide()
heading(s9, "Results: Evolved vs All Baselines",
        "Experiment: 260312_50i_openevolve_v2wl  •  OpenEvolve  •  2 seeds  •  2 workloads")

# Headline comparison: evolved vs 1:1
for bi, (label, score, sub, col) in enumerate([
    ("1:1 Baseline\n(hand-tuned)", "0%", "Control — best before AI", RED),
    ("Evolved\n(iter 16)", "+11.5%", "Best of 50 iterations", GREEN),
]):
    bx = 0.5 + bi * 5.5
    box(s9, bx, 1.28, 5.0, 1.85, fill=CARD_BG,
        border=col, border_w=Pt(2.5 if bi==1 else 1.5))
    txt(s9, label,  bx+0.2, 1.35, 4.6, 0.52, size=14, color=LIGHT)
    txt(s9, score,  bx+0.2, 1.82, 4.6, 0.75, size=40, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s9, sub,    bx+0.2, 2.60, 4.6, 0.36, size=12, color=LIGHT, align=PP_ALIGN.CENTER)

txt(s9, "→", 5.55, 1.85, 0.9, 0.7, size=36, color=ACCENT, align=PP_ALIGN.CENTER)
txt(s9, "+11.5%", 5.42, 2.58, 1.2, 0.4, size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# Multi-baseline table
txt(s9, "Full Baseline Comparison  (lower E2E = better)", 0.35, 3.22, 12.5, 0.38,
    size=14, bold=True, color=ACCENT)

# Headers
col_hdrs = ["Baseline", "glia_40qps E2E", "glia_40qps P95", "prefix_heavy E2E", "prefix_heavy P95", "vs 1:1"]
col_xs   = [0.35, 2.30, 4.10, 5.90, 8.05, 10.55]
col_ws   = [1.90, 1.75, 1.75, 2.10, 2.45, 1.70]

for ci, (hdr, cx, cw) in enumerate(zip(col_hdrs, col_xs, col_ws)):
    box(s9, cx, 3.68, cw-0.06, 0.36, fill=RGBColor(0x10,0x25,0x3a), border=ACCENT, border_w=Pt(0.8))
    txt(s9, hdr, cx+0.06, 3.70, cw-0.12, 0.30, size=11, bold=True,
        color=ACCENT, align=PP_ALIGN.CENTER)

# Data rows: LLQ, Glia, 3:2:2, 1:1, Evolved
rows_data = [
    ("LLQ",          "6,357 ms", "25,400 ms", "1,300 ms", "2,600 ms",  "−61.0%",  RED,   RGBColor(0x30,0x08,0x08)),
    ("Glia",         "4,457 ms", "18,000 ms", "  880 ms", "1,760 ms",  "−14.3%",  ORANGE,RGBColor(0x28,0x18,0x08)),
    ("3:2:2",        "4,311 ms", "17,300 ms", "  818 ms", "1,900 ms",  " −0.7%",  YELLOW,RGBColor(0x24,0x24,0x08)),
    ("1:1  ★control","4,314 ms", "17,241 ms", "  790 ms", "1,909 ms",  "  0.0%",  LIGHT, RGBColor(0x18,0x28,0x38)),
    ("Evolved  ✓",   "4,303 ms", "16,813 ms", "  700 ms", "1,435 ms",  "+11.5%",  GREEN, RGBColor(0x0a,0x25,0x10)),
]
for ri, (name, e2e1, p95_1, e2e2, p95_2, vs, col, bg) in enumerate(rows_data):
    vals = [name, e2e1, p95_1, e2e2, p95_2, vs]
    for ci, (val, cx, cw) in enumerate(zip(vals, col_xs, col_ws)):
        box(s9, cx, 4.09 + ri*0.48, cw-0.06, 0.44,
            fill=bg, border=RGBColor(0x33,0x55,0x66), border_w=Pt(0.4))
        vc = col if ci in (0, 5) else WHITE
        fw = (ci == 5)
        txt(s9, val, cx+0.06, 4.11+ri*0.48, cw-0.12, 0.38, size=11,
            color=vc, align=PP_ALIGN.CENTER, bold=fw)

# Callout: where gain comes from
box(s9, 0.35, 6.55, 12.6, 0.68, fill=RGBColor(0x0a,0x22,0x10), border=GREEN, border_w=Pt(1.5))
txt(s9,
    "Key insight: Gain is concentrated on prefix_heavy (+11.4% E2E, +24.8% P95). "
    "Adaptive prefix-affinity decay prevents hotspotting on cached instances. "
    "General traffic (glia_40qps) gains only +0.3% — load balance already works well there.",
    0.5, 6.57, 12.3, 0.62, size=12, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Deployment Viability
# ═══════════════════════════════════════════════════════════════════════════
s10 = slide()
heading(s10, "Does It Transfer to Real llm-d?",
        "Signal availability and deployment viability assessment")

signals = [
    ("InFlightRequests", "PRIMARY load signal",
     "Router-local counter, updated every request.\nNo Prometheus — zero staleness.\nDirect drop-in to llm-d routing plugin.",
     GREEN, "HIGH"),
    ("KVUtilization", "Memory pressure signal",
     "Prometheus metric, ~5s refresh cycle.\nUsed only for penalty when KVUtil > 90%.\nTolerant of stale values — edge case only.",
     ACCENT, "HIGH"),
    ("Prefix-affinity\nscore", "Cache locality signal",
     "Computed by router-local LRU cache index.\nFresh every request, zero Prometheus dep.\nAvailable in llm-d natively.",
     ACCENT, "HIGH"),
    ("QueueDepth\nBatchSize", "Load proxies (unused)",
     "Available via Prometheus (~5s stale).\nNot used by the best evolved policy.\nPresent in initial program for LLM context.",
     YELLOW, "N/A"),
    ("OutputTokens\nCacheHitRate per-req", "Simulator-only",
     "Not exposed in real llm-d router interface.\nBest evolved policy does NOT use these.\nSafe to ignore — no integration needed.",
     LIGHT, "N/A"),
]

for i, (sig, role, desc, col, viability) in enumerate(signals):
    bx = 0.35 + (i % 3) * 4.3
    by = 1.32 + (i // 3) * 2.55
    box(s10, bx, by, 4.15, 2.3, fill=CARD_BG, border=col, border_w=Pt(1.2))
    v_col = GREEN if viability=="HIGH" else (YELLOW if viability=="MED" else LIGHT)
    box(s10, bx + 2.7, by+0.08, 1.3, 0.38, fill=v_col, border=v_col)
    txt(s10, viability, bx+2.7, by+0.09, 1.3, 0.34, size=12, bold=True,
        color=BG, align=PP_ALIGN.CENTER)
    txt(s10, sig,  bx+0.15, by+0.08, 2.5, 0.52, size=12, bold=True, color=col)
    txt(s10, role, bx+0.15, by+0.60, 3.9, 0.28, size=11, italic=True, color=WHITE)
    txt(s10, desc, bx+0.15, by+0.92, 3.9, 0.88, size=11, color=LIGHT)

box(s10, 0.35, 6.62, 12.6, 0.6, fill=RGBColor(0x0a,0x25,0x10), border=GREEN, border_w=Pt(1.5))
txt(s10,
    "✅  Deployment verdict: HIGH viability.  The best policy uses only InFlightRequests (router-local, fresh) "
    "+ KVUtilization (Prometheus, stale-tolerant) + prefix scores (router-local, fresh). "
    "No simulator-only signals. Drop-in replacement for WeightedScoring.Route().",
    0.5, 6.65, 12.3, 0.52, size=12, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Summary
# ═══════════════════════════════════════════════════════════════════════════
s11 = slide()
heading(s11, "Summary", "What we set up, what was discovered, what it means")

takeaways = [
    (GREEN,  "Setup",
     "4-instance qwen2.5-7b cluster with two workloads: general bursty traffic (40 QPS) "
     "and prefix-heavy traffic (85 QPS, 6 shared-prefix groups). "
     "Score: mean of E2E and P95 latency improvement vs 1:1 baseline. Two seeds."),
    (ACCENT, "Discovery",
     "OpenEvolve found an adaptive routing policy in 50 iterations (32 min). "
     "Core insight: decay prefix-affinity weight when the cached instance is overloaded, "
     "using InFlightRequests as a fresh load signal. Plus KV pressure penalty."),
    (YELLOW, "Gains",
     "+11.5% vs best hand-tuned 1:1 baseline. +17.8% vs Glia. +42.0% vs LLQ (common default). "
     "Gains concentrated on prefix_heavy (+11.4% E2E, +24.8% P95). "
     "General traffic improves only slightly — load balance already works there."),
    (ORANGE, "Why it works",
     "Naive 1:1 routing can hotspot the instance with best cache, causing queuing. "
     "The evolved policy detects this via InFlightRequests delta and redistributes load. "
     "KV penalty kicks in only at extreme memory pressure (>90%) — a safety valve."),
    (RGBColor(0xDD,0x88,0xFF), "Next steps",
     "Validate on real llm-d with actual prefix cache hits. "
     "Test with more prefix groups and higher request rates. "
     "Explore joint router + admission co-evolution. "
     "Run with more iterations to see if >11.5% is achievable."),
]

for i, (col, label, text) in enumerate(takeaways):
    by = 1.32 + i * 1.12
    box(s11, 0.35, by, 0.72, 0.88, fill=col, border=col)
    txt(s11, label, 0.38, by+0.18, 0.68, 0.52, size=11, bold=True,
        color=BG, align=PP_ALIGN.CENTER)
    box(s11, 1.12, by, 11.85, 0.88, fill=CARD_BG, border=col, border_w=Pt(1))
    txt(s11, text, 1.27, by+0.12, 11.55, 0.65, size=12.5, color=LIGHT)


# ── Save ─────────────────────────────────────────────────────────────────────
out = "docs/slides/blis_router_experiment.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
